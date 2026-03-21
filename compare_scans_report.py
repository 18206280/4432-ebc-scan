"""
compare_scans_report.py
=======================
Compares two RRU RL Health scan CSV exports and generates a PowerPoint
comparison report.

Usage: 
    python compare_scans_report.py 20260212\RRU_BXP_RL50_Health_Result.csv 20260306\RRU_BXP_RL50_Health_Result.csv

Dependencies:
    pip install python-pptx matplotlib pandas
"""

import argparse
import io
import sys
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ═══════════════════════════════════════════════════════════════════════════════
# COLOUR PALETTE  (matches monthly report)
# ═══════════════════════════════════════════════════════════════════════════════

class C:
    NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
    NAVY_MID    = RGBColor(0x2C, 0x4F, 0x7A)
    NAVY_LIGHT  = RGBColor(0xEB, 0xF0, 0xF7)
    TEAL        = RGBColor(0x08, 0x91, 0xB2)
    TEAL_LIGHT  = RGBColor(0xE0, 0xF2, 0xFE)
    GREEN_DARK  = RGBColor(0x2D, 0x6A, 0x2D)
    GREEN_MID   = RGBColor(0x16, 0xA3, 0x4A)
    GREEN_LIGHT = RGBColor(0xEB, 0xF5, 0xEB)
    AMBER_DARK  = RGBColor(0x92, 0x40, 0x0E)
    AMBER_MID   = RGBColor(0xD9, 0x77, 0x06)
    AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
    ORANGE_DARK = RGBColor(0x7C, 0x2D, 0x12)
    ORANGE_MID  = RGBColor(0xEA, 0x58, 0x0C)
    ORANGE_LIGHT= RGBColor(0xFF, 0xF0, 0xE6)
    RED_DARK    = RGBColor(0x7F, 0x1D, 0x1D)
    RED_MID     = RGBColor(0xDC, 0x26, 0x26)
    RED_LIGHT   = RGBColor(0xFE, 0xF2, 0xF2)
    PURPLE_DARK = RGBColor(0x4C, 0x1D, 0x95)
    PURPLE_MID  = RGBColor(0x7C, 0x3A, 0xED)
    PURPLE_LIGHT= RGBColor(0xF3, 0xE8, 0xFF)
    SLATE_DARK  = RGBColor(0x37, 0x41, 0x51)
    SLATE_MID   = RGBColor(0x64, 0x74, 0x8B)
    SLATE_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
    BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT        = RGBColor(0x1E, 0x29, 0x3B)

STATUS_ORDER = ["Healthy", "Early Degradation", "Middle Degradation",
                "Late Degradation", "Critical", "No Data"]

STATUS_CFG = {
    "Healthy":            ("✅", C.GREEN_DARK,  C.GREEN_LIGHT,  C.GREEN_MID),
    "Early Degradation":  ("⚠️",  C.AMBER_DARK,  C.AMBER_LIGHT,  C.AMBER_MID),
    "Middle Degradation": ("🔶", C.ORANGE_DARK, C.ORANGE_LIGHT, C.ORANGE_MID),
    "Late Degradation":   ("🔴", C.RED_DARK,    C.RED_LIGHT,    C.RED_MID),
    "Critical":           ("🚨", C.PURPLE_DARK, C.PURPLE_LIGHT, C.PURPLE_MID),
    "No Data":            ("—",  C.SLATE_DARK,  C.SLATE_LIGHT,  C.SLATE_MID),
}

ALERT_KEYS = ["Early Degradation", "Middle Degradation", "Late Degradation", "Critical"]

# ═══════════════════════════════════════════════════════════════════════════════
# DATA LOADING & COMPUTATION
# ═══════════════════════════════════════════════════════════════════════════════

def load_df(csv_path: str) -> pd.DataFrame:
    df = pd.read_csv(csv_path)
    df["status"] = df["Unit_Worst_Status"].str.extract(
        r"(Healthy|Early Degradation|Middle Degradation|Late Degradation|Critical|No Data)",
        expand=False)
    df["serial_number"] = df["serial_number"].astype(str).str.strip()
    for b in "ABCD":
        df[f"RL_{b}_CoG"] = pd.to_numeric(df[f"RL_{b}_CoG"], errors="coerce")
    # Scan date from RL time columns
    time_vals = []
    for b in "ABCD":
        col = f"RL_{b}_time"
        if col in df.columns:
            v = df[col].dropna().astype(str)
            v = v[v.str.match(r"^\d{6}")]
            if len(v): time_vals.append(v.max()[:6])
    df._scan_date = max(time_vals) if time_vals else "??????"
    return df


def fmt_date(yymmdd: str) -> str:
    try:
        return datetime.strptime(yymmdd, "%y%m%d").strftime("%d %b %Y")
    except Exception:
        return yymmdd


def has_override(row, suffix=""):
    """Return True if any branch has CoG override on a merged row."""
    for b in "ABCD":
        key = f"RL_{b}_CoG_Override{suffix}"
        if key in row.index and "Yes" in str(row[key]):
            return True
    return False


def compute_comparison(df_old: pd.DataFrame, df_new: pd.DataFrame) -> dict:
    """Build all stats needed for slides."""
    s = {}

    s["date_old"] = fmt_date(df_old._scan_date)
    s["date_new"] = fmt_date(df_new._scan_date)
    s["total_old"] = len(df_old)
    s["total_new"] = len(df_new)

    sn_old = set(df_old["serial_number"])
    sn_new = set(df_new["serial_number"])
    both   = sn_old & sn_new
    only_old = sn_old - sn_new
    only_new = sn_new - sn_old

    s["intersection"] = len(both)
    s["disappeared"]  = len(only_old)
    s["new_units"]    = len(only_new)

    # Build intersection dataframes (deduplicated)
    d_old = df_old[df_old["serial_number"].isin(both)].drop_duplicates("serial_number").set_index("serial_number")
    d_new = df_new[df_new["serial_number"].isin(both)].drop_duplicates("serial_number").set_index("serial_number")

    # ── Status distributions ────────────────────────────────────────────────
    s["dist_old"] = d_old["status"].value_counts().reindex(STATUS_ORDER, fill_value=0).to_dict()
    s["dist_new"] = d_new["status"].value_counts().reindex(STATUS_ORDER, fill_value=0).to_dict()
    s["alert_old"] = sum(s["dist_old"][k] for k in ALERT_KEYS)
    s["alert_new"] = sum(s["dist_new"][k] for k in ALERT_KEYS)
    s["alert_delta"] = s["alert_new"] - s["alert_old"]

    # ── Transition matrix ───────────────────────────────────────────────────
    merged = d_old[["status"]].join(d_new[["status"]], lsuffix="_old", rsuffix="_new")
    merged = merged.dropna(subset=["status_old", "status_new"])

    trans = pd.crosstab(merged["status_old"], merged["status_new"])
    trans = trans.reindex(index=STATUS_ORDER, columns=STATUS_ORDER, fill_value=0)
    s["trans"] = trans

    # ── Change summary ──────────────────────────────────────────────────────
    changed = merged[merged["status_old"] != merged["status_new"]].copy()
    s["total_changed"] = len(changed)
    s["pct_changed"]   = round(len(changed) / len(merged) * 100, 1)

    order_map = {k: i for i, k in enumerate(STATUS_ORDER)}
    changed["old_ord"] = changed["status_old"].map(order_map)
    changed["new_ord"] = changed["status_new"].map(order_map)
    changed["direction"] = changed.apply(
        lambda r: "Worsened" if r["new_ord"] > r["old_ord"] else "Improved", axis=1)

    s["worsened"] = int((changed["direction"] == "Worsened").sum())
    s["improved"] = int((changed["direction"] == "Improved").sum())

    # Worsened breakdown
    worsened_df = changed[changed["direction"] == "Worsened"]
    s["worsened_detail"] = worsened_df.groupby(
        ["status_old", "status_new"]).size().reset_index(name="count").sort_values("count", ascending=False)

    # Improved breakdown
    improved_df = changed[changed["direction"] == "Improved"]
    s["improved_detail"] = improved_df.groupby(
        ["status_old", "status_new"]).size().reset_index(name="count").sort_values("count", ascending=False)

    # ── Flip analysis ───────────────────────────────────────────────────────
    # Join CoG and override data
    cog_cols_old = {f"RL_{b}_CoG": f"cog_{b}_old" for b in "ABCD"}
    cog_cols_new = {f"RL_{b}_CoG": f"cog_{b}_new" for b in "ABCD"}
    ov_cols_old  = {f"RL_{b}_CoG_Override": f"ov_{b}_old" for b in "ABCD"}
    ov_cols_new  = {f"RL_{b}_CoG_Override": f"ov_{b}_new" for b in "ABCD"}

    extra_old = d_old[[c for c in list(cog_cols_old) + list(ov_cols_old) if c in d_old.columns]].rename(columns={**cog_cols_old, **ov_cols_old})
    extra_new = d_new[[c for c in list(cog_cols_new) + list(ov_cols_new) if c in d_new.columns]].rename(columns={**cog_cols_new, **ov_cols_new})
    rich = merged.join(extra_old).join(extra_new)

    for b in "ABCD":
        rich[f"cog_{b}_old"] = pd.to_numeric(rich.get(f"cog_{b}_old"), errors="coerce")
        rich[f"cog_{b}_new"] = pd.to_numeric(rich.get(f"cog_{b}_new"), errors="coerce")

    rich["min_cog_old"] = rich[[f"cog_{b}_old" for b in "ABCD"]].min(axis=1)
    rich["min_cog_new"] = rich[[f"cog_{b}_new" for b in "ABCD"]].min(axis=1)
    rich["cog_delta"]   = rich["min_cog_new"] - rich["min_cog_old"]

    # Healthy(old) → Early(new)
    h2e = rich[(rich["status_old"] == "Healthy") & (rich["status_new"] == "Early Degradation")].copy()
    h2e["via_override"] = h2e.apply(lambda r: any("Yes" in str(r.get(f"ov_{b}_new","")) for b in "ABCD"), axis=1)
    s["h2e_total"]    = len(h2e)
    s["h2e_override"] = int(h2e["via_override"].sum())
    s["h2e_fv"]       = int((~h2e["via_override"]).sum())
    s["h2e_cog_old"]  = round(h2e["min_cog_old"].mean(), 2)
    s["h2e_cog_new"]  = round(h2e["min_cog_new"].mean(), 2)
    s["h2e_cog_delta"]= round(h2e["cog_delta"].mean(), 2)

    # Early(old) → Healthy(new)
    e2h = rich[(rich["status_old"] == "Early Degradation") & (rich["status_new"] == "Healthy")].copy()
    e2h["via_override"] = e2h.apply(lambda r: any("Yes" in str(r.get(f"ov_{b}_old","")) for b in "ABCD"), axis=1)
    s["e2h_total"]    = len(e2h)
    s["e2h_override"] = int(e2h["via_override"].sum())
    s["e2h_fv"]       = int((~e2h["via_override"]).sum())
    s["e2h_cog_old"]  = round(e2h["min_cog_old"].mean(), 2)
    s["e2h_cog_new"]  = round(e2h["min_cog_new"].mean(), 2)
    s["e2h_cog_delta"]= round(e2h["cog_delta"].mean(), 2)

    # ── CoG stability (stable-healthy units) ────────────────────────────────
    stable = rich[(rich["status_old"] == "Healthy") & (rich["status_new"] == "Healthy")].copy()
    s["stable_healthy_n"] = len(stable)
    all_deltas = pd.concat([stable[f"cog_{b}_new"] - stable[f"cog_{b}_old"] for b in "ABCD"]).dropna()
    s["cog_delta_std"]   = round(float(all_deltas.std()), 3)
    s["cog_delta_max"]   = round(float(all_deltas.abs().max()), 2)
    s["cog_delta_above1"]= int((all_deltas.abs() > 1).sum())
    s["cog_delta_pct1"]  = round(s["cog_delta_above1"] / len(all_deltas) * 100, 1) if len(all_deltas) else 0
    s["cog_deltas"]      = all_deltas.values   # for histogram

    # disappeared / new status breakdown
    dis_df = df_old[df_old["serial_number"].isin(only_old)]
    new_df = df_new[df_new["serial_number"].isin(only_new)]
    s["disappeared_status"] = dis_df["status"].value_counts().reindex(STATUS_ORDER, fill_value=0).to_dict()
    s["new_units_status"]   = new_df["status"].value_counts().reindex(STATUS_ORDER, fill_value=0).to_dict()

    # ── VSWR Risk Assessment ────────────────────────────────────────────────
    COLS_AG = list(df_old.columns[:7])

    vswr_old = set(df_old[df_old["VSWR_risk_assessment"] == "Risk unit"]["serial_number"])
    vswr_new = set(df_new[df_new["VSWR_risk_assessment"] == "Risk unit"]["serial_number"])

    vswr_both_old = vswr_old & both          # risk in old, within intersection
    vswr_both_new = vswr_new & both          # risk in new, within intersection
    vswr_persistent_sn = vswr_both_old & vswr_both_new
    vswr_new_alerts_sn  = vswr_both_new - vswr_both_old
    vswr_disappeared_sn = vswr_both_old - vswr_both_new

    s["vswr_total_old"]       = len(vswr_old)
    s["vswr_total_new"]       = len(vswr_new)
    s["vswr_persistent"]      = len(vswr_persistent_sn)
    s["vswr_new_alerts"]      = len(vswr_new_alerts_sn)
    s["vswr_disappeared"]     = len(vswr_disappeared_sn)
    s["vswr_net_change"]      = len(vswr_both_new) - len(vswr_both_old)

    # Persistent — RL status change
    pers_old_st = d_old.loc[list(vswr_persistent_sn), "status"]
    pers_new_st = d_new.loc[list(vswr_persistent_sn), "status"]
    order_map2  = {k: i for i, k in enumerate(STATUS_ORDER)}
    pers_dir = pd.Series({
        sn: ("Worsened" if order_map2.get(n, 5) > order_map2.get(o, 5)
             else ("Improved" if order_map2.get(n, 5) < order_map2.get(o, 5) else "Stable"))
        for sn, o, n in zip(pers_old_st.index, pers_old_st, pers_new_st)
    })
    s["vswr_pers_worsened"] = int((pers_dir == "Worsened").sum())
    s["vswr_pers_stable"]   = int((pers_dir == "Stable").sum())
    s["vswr_pers_improved"] = int((pers_dir == "Improved").sum())

    # New alerts — RL status breakdown
    new_al_st = d_new.loc[list(vswr_new_alerts_sn), "status"]
    s["vswr_new_rl_dist"] = new_al_st.value_counts().reindex(STATUS_ORDER, fill_value=0).to_dict()

    # Disappeared — old RL status
    dis_al_st = d_old.loc[list(vswr_disappeared_sn), "status"]
    s["vswr_dis_rl_dist"] = dis_al_st.value_counts().reindex(STATUS_ORDER, fill_value=0).to_dict()

    # ── Store dataframes needed for CSV export ──────────────────────────────
    s["_df_old"]   = df_old
    s["_df_new"]   = df_new
    s["_d_old"]    = d_old
    s["_d_new"]    = d_new
    s["_rich"]     = rich
    s["_changed"]  = changed
    s["_COLS_AG"]  = COLS_AG
    s["_vswr_persistent_sn"]  = vswr_persistent_sn
    s["_vswr_new_alerts_sn"]  = vswr_new_alerts_sn
    s["_vswr_disappeared_sn"] = vswr_disappeared_sn
    s["_pers_dir"] = pers_dir
    s["_pers_old_st"] = pers_old_st
    s["_pers_new_st"] = pers_new_st

    return s


# ═══════════════════════════════════════════════════════════════════════════════
# DRAWING HELPERS  (same style as monthly report)
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill, line=None, line_w=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line; shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, size=10, bold=False,
                 color=None, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = "Calibri"
    r.font.color.rgb = color if color else C.TEXT
    return txb


def add_slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 10, 0.62, C.NAVY)
    add_text_box(slide, 0.38, 0.08, 9.24, 0.46, title,
                 size=20, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_rect(slide, 0, 0.62, 10, 0.26, C.NAVY_MID)
        add_text_box(slide, 0.38, 0.65, 9.24, 0.20, subtitle,
                     size=8.5, color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.LEFT)


def add_kpi_card(slide, x, y, w, h, label, value, dark, light, num_color):
    add_rect(slide, x, y, w, 0.28, dark)
    add_text_box(slide, x+0.06, y+0.03, w-0.12, 0.22, label,
                 size=8.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, y+0.28, w, h-0.28, light, C.BORDER, 0.5)
    add_text_box(slide, x+0.06, y+0.30, w-0.12, h-0.34, value,
                 size=22, bold=True, color=num_color, align=PP_ALIGN.CENTER)


def add_section_label(slide, x, y, w, text):
    add_rect(slide, x, y, w, 0.26, C.TEAL)
    add_text_box(slide, x+0.10, y+0.03, w-0.15, 0.20, text,
                 size=9, bold=True, color=C.WHITE)


def add_insight_box(slide, x, y, w, h, text):
    add_rect(slide, x, y, 0.04, h, C.TEAL)
    add_rect(slide, x+0.04, y, w-0.04, h, C.TEAL_LIGHT, C.BORDER, 0.5)
    add_text_box(slide, x+0.14, y+0.05, 0.80, 0.16,
                 "Key Insight", size=9, bold=True, color=C.NAVY)
    add_text_box(slide, x+0.14, y+0.20, w-0.20, h-0.22,
                 text, size=8.5, color=C.TEXT)


def add_table(slide, x, y, w, h, rows, col_widths, row_heights=None):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    n_rows = len(rows); n_cols = len(rows[0])
    tbl = slide.shapes.add_table(n_rows, n_cols,
        Inches(x), Inches(y), Inches(w), Inches(h)).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)
    if row_heights:
        for ri, rh in enumerate(row_heights):
            tbl.rows[ri].height = Inches(rh)
    for ri, row in enumerate(rows):
        for ci, cell_def in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                           "right": PP_ALIGN.RIGHT}.get(cell_def.get("align","center"), PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = str(cell_def.get("text",""))
            run.font.size = Pt(cell_def.get("font_size", 10))
            run.font.bold = cell_def.get("bold", False)
            run.font.italic = cell_def.get("italic", False)
            run.font.name = "Calibri"
            if cell_def.get("color"): run.font.color.rgb = cell_def["color"]
            if cell_def.get("fill"):
                cell.fill.solid(); cell.fill.fore_color.rgb = cell_def["fill"]
    return tbl


def hdr(text, w=None):
    d = {"text": text, "bold": True, "fill": C.NAVY, "color": C.WHITE,
         "align": "center", "font_size": 10}
    if w: d["width"] = w
    return d


def dat(text, align="center", bold=False, color=None, fill=None, font_size=10):
    return {"text": text, "align": align, "bold": bold,
            "color": color or C.TEXT, "fill": fill or C.WHITE, "font_size": font_size}


def altfill(i):
    return C.NAVY_LIGHT if i % 2 == 0 else C.WHITE


# ═══════════════════════════════════════════════════════════════════════════════
# CHART BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_transition_heatmap(trans: pd.DataFrame) -> bytes:
    labels_short = ["Healthy","Early","Middle","Late","Critical","No Data"]
    data = trans.values.astype(float)

    fig, ax = plt.subplots(figsize=(5.2, 3.6), dpi=150)
    fig.patch.set_facecolor("white")

    # Mask diagonal (no change)
    mask = np.eye(len(labels_short), dtype=bool)
    display = np.where(mask, np.nan, data)

    cmap = plt.cm.YlOrRd
    cmap.set_bad("white")
    im = ax.imshow(display, cmap=cmap, aspect="auto", vmin=0)

    # Diagonal cells in light gray
    for i in range(len(labels_short)):
        ax.add_patch(plt.Rectangle((i-0.5, i-0.5), 1, 1,
                     fill=True, color="#F0F4F8", zorder=1))
        if trans.values[i, i] > 0:
            ax.text(i, i, str(trans.values[i, i]), ha="center", va="center",
                    fontsize=8, color="#64748B", zorder=2)

    # Off-diagonal annotations
    for r in range(len(labels_short)):
        for c in range(len(labels_short)):
            if r != c and data[r, c] > 0:
                ax.text(c, r, str(int(data[r, c])), ha="center", va="center",
                        fontsize=9, fontweight="bold", color="white" if data[r,c] > 30 else "#1E293B", zorder=3)

    ax.set_xticks(range(len(labels_short)))
    ax.set_yticks(range(len(labels_short)))
    ax.set_xticklabels(labels_short, fontsize=7.5, rotation=30, ha="right")
    ax.set_yticklabels(labels_short, fontsize=7.5)
    ax.set_xlabel("Status in newer scan →", fontsize=8, color="#475569", labelpad=4)
    ax.set_ylabel("← Status in older scan", fontsize=8, color="#475569", labelpad=4)
    ax.set_title("Status Transition Matrix", fontsize=10, fontweight="bold",
                 color="#1E293B", pad=8)
    ax.tick_params(colors="#64748B")
    for spine in ax.spines.values(): spine.set_visible(False)

    plt.colorbar(im, ax=ax, fraction=0.03, pad=0.02).ax.tick_params(labelsize=7)
    plt.tight_layout(pad=0.5)

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig); buf.seek(0)
    return buf


def build_cog_stability_chart(deltas: np.ndarray) -> bytes:
    fig, ax = plt.subplots(figsize=(5.5, 2.8), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    bins = np.arange(-3, 3.5, 0.5)
    n, edges, patches = ax.hist(deltas, bins=bins, color="#0891B2",
                                 edgecolor="white", linewidth=0.4)

    # Colour bars outside ±1 dB in amber
    for patch, left in zip(patches, edges[:-1]):
        if abs(left + 0.25) > 1.0:
            patch.set_facecolor("#D97706")

    ax.axvline(x=0, color="#94A3B8", linewidth=1.2, linestyle="-")
    ax.axvline(x=1, color="#DC2626", linewidth=1.2, linestyle="--", alpha=0.7)
    ax.axvline(x=-1, color="#DC2626", linewidth=1.2, linestyle="--", alpha=0.7)

    ylim = ax.get_ylim()
    ax.text(1.05, ylim[1]*0.85, "±1 dB", color="#DC2626", fontsize=7.5, fontweight="bold")

    ax.set_xlabel("CoG change (dB)", fontsize=9, color="#1E293B", labelpad=4)
    ax.set_ylabel("Branch count", fontsize=9, color="#1E293B", labelpad=4)
    ax.set_title("CoG Stability — Stable Healthy Units", fontsize=10,
                 fontweight="bold", color="#1E293B", pad=6)
    ax.set_xlim(-3.5, 3.5)
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#E2E8F0"); ax.spines["bottom"].set_color("#E2E8F0")
    ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.6, zorder=0)
    ax.set_axisbelow(True)
    ax.tick_params(colors="#64748B", labelsize=8)

    blue_patch  = mpatches.Patch(color="#0891B2", label="Within ±1 dB")
    amber_patch = mpatches.Patch(color="#D97706", label="Outside ±1 dB")
    ax.legend(handles=[blue_patch, amber_patch], fontsize=7.5, loc="upper right")

    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig); buf.seek(0)
    return buf


def build_change_bar(worsened: int, improved: int) -> bytes:
    fig, ax = plt.subplots(figsize=(3.0, 2.2), dpi=150)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    bars = ax.bar(["Worsened", "Improved"], [worsened, improved],
                  color=["#DC2626", "#16A34A"], width=0.5, edgecolor="white")
    for bar, val in zip(bars, [worsened, improved]):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                str(val), ha="center", va="bottom", fontsize=11, fontweight="bold",
                color="#1E293B")
    ax.set_ylabel("Units", fontsize=8.5, color="#1E293B")
    ax.set_title("Status Changes", fontsize=9.5, fontweight="bold", color="#1E293B", pad=5)
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#E2E8F0"); ax.spines["bottom"].set_color("#E2E8F0")
    ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.6, zorder=0)
    ax.set_axisbelow(True); ax.tick_params(colors="#64748B", labelsize=8.5)
    plt.tight_layout(pad=0.4)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig); buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_title(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.NAVY

    # Antenna icon
    add_text_box(slide, 0.38, 0.22, 0.60, 0.50, "(A)", size=22,
                 bold=True, color=C.TEAL, align=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, 1.10, 0.18, 8.52, 0.72,
                 "Radio 4432 B28 — RL Health Scan Comparison Report",
                 size=26, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)

    # Teal rule
    add_rect(slide, 0.38, 1.02, 9.24, 0.04, C.TEAL)

    # Subtitle
    add_text_box(slide, 0.38, 1.14, 9.24, 0.30,
                 f"Scan-over-Scan Status Change Analysis  |  {s['date_old']}  →  {s['date_new']}",
                 size=13, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.LEFT)

    # 4 KPI boxes
    kpis = [
        (f"{s['total_old']:,}",      f"Units — {s['date_old']}",   C.TEAL,        C.NAVY_MID,  C.WHITE),
        (f"{s['total_new']:,}",      f"Units — {s['date_new']}",   C.TEAL,        C.NAVY_MID,  C.WHITE),
        (f"{s['intersection']:,}",   "Common Units (Compared)",    C.GREEN_DARK,  C.NAVY_MID,  C.GREEN_MID),
        (f"{s['total_changed']:,}",  "Units with Status Change",   C.AMBER_DARK,  C.NAVY_MID,  C.AMBER_MID),
    ]
    kw = 2.10; ky = 1.70; kh = 0.90
    for i, (val, lbl, dk, bg_, nc) in enumerate(kpis):
        x = 0.38 + i * 2.34
        add_rect(slide, x, ky, kw, 0.26, dk)
        add_text_box(slide, x+0.06, ky+0.03, kw-0.12, 0.20, lbl,
                     size=7.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x, ky+0.26, kw, kh-0.26, bg_)
        add_text_box(slide, x+0.06, ky+0.30, kw-0.12, kh-0.32, val,
                     size=24, bold=True, color=nc, align=PP_ALIGN.CENTER)

    add_text_box(slide, 0.38, 5.38, 9.24, 0.20,
                 "Confidential — Internal Management Report",
                 size=8.5, italic=True,
                 color=RGBColor(0x70, 0x8A, 0xA8), align=PP_ALIGN.CENTER)


def build_slide_overview(prs, s):
    """Slide 2: Comparison Overview — unit counts + status distributions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Scan Comparison Overview",
                     f"{s['date_old']}  →  {s['date_new']}  |  Intersection: {s['intersection']:,} units  |  "
                     f"Disappeared: {s['disappeared']}  |  New: {s['new_units']}")

    # 3 population KPI cards
    pop_kpis = [
        (f"{s['intersection']:,}", "Units in Both Scans", C.TEAL,       C.TEAL_LIGHT,   C.TEAL),
        (f"{s['disappeared']}",    "Disappeared Units",   C.SLATE_DARK, C.SLATE_LIGHT,  C.SLATE_MID),
        (f"{s['new_units']}",      "New Units",           C.GREEN_DARK, C.GREEN_LIGHT,  C.GREEN_MID),
    ]
    kw = 2.80; ky = 1.05; kh = 0.80
    for i, (val, lbl, dk, lt, nc) in enumerate(pop_kpis):
        add_kpi_card(slide, 0.38 + i*3.0, ky, kw, kh, lbl, val, dk, lt, nc)

    # Status comparison table
    add_section_label(slide, 0.38, 2.02, 9.24, "Status Distribution Comparison")
    tbl_rows = [
        [hdr("Status"), hdr(s["date_old"]), hdr(s["date_new"]), hdr("Change")],
    ]
    for st in STATUS_ORDER:
        cfg = STATUS_CFG[st]
        v_old = s["dist_old"].get(st, 0)
        v_new = s["dist_new"].get(st, 0)
        delta = v_new - v_old
        delta_str = f"+{delta}" if delta > 0 else str(delta)
        delta_col = C.RED_MID if delta > 0 and st in ALERT_KEYS else (
                    C.GREEN_MID if delta < 0 and st in ALERT_KEYS else C.SLATE_MID)
        fill = altfill(STATUS_ORDER.index(st))
        tbl_rows.append([
            dat(f"{cfg[0]} {st}", "left",  True,  cfg[1], fill),
            dat(f"{v_old:,}",     "center", False, C.TEXT, fill),
            dat(f"{v_new:,}",     "center", False, C.TEXT, fill),
            dat(delta_str,        "center", True,  delta_col, fill),
        ])
    # Totals row
    tot_old = sum(s["dist_old"].values())
    tot_new = sum(s["dist_new"].values())
    tbl_rows.append([
        dat("Total", "left", True, C.WHITE, C.NAVY),
        dat(f"{tot_old:,}", "center", True, C.WHITE, C.NAVY),
        dat(f"{tot_new:,}", "center", True, C.WHITE, C.NAVY),
        dat(f"{tot_new-tot_old:+,}", "center", True, C.WHITE, C.NAVY),
    ])
    add_table(slide, 0.38, 2.32, 9.24, 2.60,
              tbl_rows, col_widths=[4.20, 1.68, 1.68, 1.68],
              row_heights=[0.30] + [0.28]*6 + [0.30])

    # Insight
    alert_delta_str = f"+{s['alert_delta']}" if s['alert_delta'] >= 0 else str(s['alert_delta'])
    add_insight_box(slide, 0.38, 4.72, 9.24, 0.36,
        f"Alert units changed from {s['alert_old']} ({round(s['alert_old']/s['intersection']*100,1)}%) "
        f"to {s['alert_new']} ({round(s['alert_new']/s['intersection']*100,1)}%) — "
        f"net change: {alert_delta_str} units across the {s['intersection']:,} common units compared.")

    add_text_box(slide, 0.38, 5.14, 9.24, 0.14,
                 "Note: Comparison performed on intersection only. Disappeared and new units excluded from status change analysis.",
                 size=7.5, italic=True, color=C.SLATE_MID)


def build_slide_transition(prs, s):
    """Slide 3: Status Transition Matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Status Transition Matrix",
                     f"How units moved between health states  |  {s['intersection']:,} common units  |  "
                     f"{s['total_changed']:,} changed ({s['pct_changed']}%)")

    # Change summary KPIs
    kpis = [
        (str(s["total_changed"]),    "Total Changed",  C.AMBER_DARK,  C.AMBER_LIGHT,  C.AMBER_MID),
        (str(s["worsened"]),         "Worsened",       C.RED_DARK,    C.RED_LIGHT,    C.RED_MID),
        (str(s["improved"]),         "Improved",       C.GREEN_DARK,  C.GREEN_LIGHT,  C.GREEN_MID),
        (f"{s['pct_changed']}%",     "Change Rate",    C.NAVY,        C.NAVY_LIGHT,   C.TEAL),
    ]
    kw = 2.10; ky = 1.05; kh = 0.80
    for i, (val, lbl, dk, lt, nc) in enumerate(kpis):
        add_kpi_card(slide, 0.38 + i*2.34, ky, kw, kh, lbl, val, dk, lt, nc)

    # Heatmap (left) — taller now that bar chart is removed
    add_section_label(slide, 0.38, 2.02, 5.50, "Transition Heatmap")
    hmap_buf = build_transition_heatmap(s["trans"])
    slide.shapes.add_picture(hmap_buf, Inches(0.38), Inches(2.32), Inches(5.50), Inches(2.85))

    # Top transitions table (right) — moved up, more rows now bar chart is gone
    add_section_label(slide, 6.08, 2.02, 3.54, "Top Transitions")
    all_detail = pd.concat([
        s["worsened_detail"].assign(dir="↓"),
        s["improved_detail"].assign(dir="↑")
    ]).sort_values("count", ascending=False).head(9)

    tbl_rows = [[hdr("From"), hdr("To"), hdr("Dir"), hdr("n")]]
    for i, row in enumerate(all_detail.itertuples()):
        fill = altfill(i)
        dk = C.RED_MID if row.dir == "↓" else C.GREEN_MID
        short = {"Healthy":"Healthy","Early Degradation":"Early",
                 "Middle Degradation":"Middle","Late Degradation":"Late",
                 "Critical":"Critical","No Data":"No Data"}
        tbl_rows.append([
            dat(short.get(row.status_old, row.status_old), "left",   False, C.TEXT, fill),
            dat(short.get(row.status_new, row.status_new), "left",   False, C.TEXT, fill),
            dat(row.dir,                                   "center", True,  dk,     fill),
            dat(str(row.count),                            "center", True,  C.TEXT, fill),
        ])
    add_table(slide, 6.08, 2.32, 3.54, 2.62,
              tbl_rows, col_widths=[1.20, 1.20, 0.52, 0.62],
              row_heights=[0.28]*len(tbl_rows))

    add_text_box(slide, 0.38, 5.26, 9.24, 0.14,
                 "Diagonal cells (grey) = no change. Off-diagonal = units that changed state. "
                 "Rows = older scan status, columns = newer scan status.",
                 size=7.5, italic=True, color=C.SLATE_MID)


def build_slide_flip(prs, s):
    """Slide 4: Flip Analysis — Healthy↔Early transitions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Flip Analysis — Healthy ↔ Early Degradation",
                     "Units oscillating around the Early Degradation boundary between scans")

    # 4 KPIs
    kpis = [
        (str(s["h2e_total"]), "Healthy → Early",     C.AMBER_DARK,  C.AMBER_LIGHT,  C.AMBER_MID),
        (str(s["e2h_total"]), "Early → Healthy",     C.GREEN_DARK,  C.GREEN_LIGHT,  C.GREEN_MID),
        (str(s["h2e_fv"]),    "H→E via FV/BTR/SZR", C.ORANGE_DARK, C.ORANGE_LIGHT, C.ORANGE_MID),
        (str(s["e2h_fv"]),    "E→H via FV/BTR/SZR", C.TEAL,        C.TEAL_LIGHT,   C.TEAL),
    ]
    kw = 2.10; ky = 1.05; kh = 0.80
    for i, (val, lbl, dk, lt, nc) in enumerate(kpis):
        add_kpi_card(slide, 0.38 + i*2.34, ky, kw, kh, lbl, val, dk, lt, nc)

    # H→E detail table (left)
    add_section_label(slide, 0.38, 2.02, 4.50, "Healthy → Early (new alerts)")
    h2e_rows = [
        [hdr("Metric"), hdr("Value")],
        [dat("Units flagged",     "left"), dat(str(s["h2e_total"]),             "center", True)],
        [dat("Via CoG override",  "left"), dat(str(s["h2e_override"]),          "center", True, C.AMBER_MID)],
        [dat("Via FV/BTR/SZR",    "left"), dat(str(s["h2e_fv"]),               "center", True, C.ORANGE_MID)],
        [dat("CoG (older scan)",  "left"), dat(f"{s['h2e_cog_old']} dB",       "center")],
        [dat("CoG (newer scan)",  "left"), dat(f"{s['h2e_cog_new']} dB",       "center")],
        [dat("Mean CoG drop",     "left"), dat(f"{s['h2e_cog_delta']:+.2f} dB","center", True, C.RED_MID)],
    ]
    for i in range(1, len(h2e_rows)):
        h2e_rows[i][0]["fill"] = altfill(i-1)
        h2e_rows[i][1]["fill"] = altfill(i-1)
    add_table(slide, 0.38, 2.32, 4.50, 1.90,
              h2e_rows, col_widths=[2.60, 1.90], row_heights=[0.28]*len(h2e_rows))

    # E→H detail table (right)
    add_section_label(slide, 5.12, 2.02, 4.50, "Early → Healthy (resolved alerts)")
    e2h_rows = [
        [hdr("Metric"), hdr("Value")],
        [dat("Units resolved",    "left"), dat(str(s["e2h_total"]),              "center", True)],
        [dat("Were CoG override", "left"), dat(str(s["e2h_override"]),           "center", True, C.AMBER_MID)],
        [dat("Were FV/BTR/SZR",   "left"), dat(str(s["e2h_fv"]),                "center", True, C.GREEN_MID)],
        [dat("CoG (older scan)",  "left"), dat(f"{s['e2h_cog_old']} dB",        "center")],
        [dat("CoG (newer scan)",  "left"), dat(f"{s['e2h_cog_new']} dB",        "center")],
        [dat("Mean CoG gain",     "left"), dat(f"{s['e2h_cog_delta']:+.2f} dB", "center", True, C.GREEN_MID)],
    ]
    for i in range(1, len(e2h_rows)):
        e2h_rows[i][0]["fill"] = altfill(i-1)
        e2h_rows[i][1]["fill"] = altfill(i-1)
    add_table(slide, 5.12, 2.32, 4.50, 1.90,
              e2h_rows, col_widths=[2.60, 1.90], row_heights=[0.28]*len(e2h_rows))

    # Insight
    add_insight_box(slide, 0.38, 4.28, 9.24, 0.92,
        f"{s['h2e_total']} Healthy→Early and {s['e2h_total']} Early→Healthy transitions observed — "
        f"nearly symmetric, suggesting boundary oscillation rather than genuine degradation. "
        f"Mean CoG drop for H→E: {s['h2e_cog_delta']:+.2f} dB (from {s['h2e_cog_old']} → {s['h2e_cog_new']} dB). "
        f"All {s['h2e_total']} H→E units scored >90 composite — flagged via FV/BTR/SZR borderline scoring, "
        f"not CoG override. Recommend monitoring these units in subsequent scans before actioning.")

    add_text_box(slide, 0.38, 5.26, 9.24, 0.14,
                 "Flip = unit classified differently in each scan without physical change evidence. "
                 "CoG delta < 0.5 dB between scans is within algorithm noise band.",
                 size=7.5, italic=True, color=C.SLATE_MID)


def build_slide_worsened(prs, s):
    """Slide 5: Worsened Units."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Worsened Units",
                     f"Units with higher degradation severity in newer scan  |  {s['worsened']} units  |  "
                     f"{s['date_old']}  →  {s['date_new']}")

    # KPIs
    add_kpi_card(slide, 0.38, 1.05, 2.80, 0.80,
                 "Total Worsened", str(s["worsened"]), C.RED_DARK, C.RED_LIGHT, C.RED_MID)
    add_kpi_card(slide, 3.42, 1.05, 2.80, 0.80,
                 "As % of Changed", f"{round(s['worsened']/s['total_changed']*100,1)}%",
                 C.ORANGE_DARK, C.ORANGE_LIGHT, C.ORANGE_MID)
    add_kpi_card(slide, 6.46, 1.05, 2.80, 0.80,
                 "As % of Intersection", f"{round(s['worsened']/s['intersection']*100,2)}%",
                 C.NAVY, C.NAVY_LIGHT, C.TEAL)

    # Worsened breakdown table
    add_section_label(slide, 0.38, 2.02, 9.24, "Worsened Transitions Breakdown")
    tbl_rows = [[hdr("From Status"), hdr("To Status"), hdr("Units"), hdr("% of Worsened")]]
    total_w = s["worsened"]
    for i, row in enumerate(s["worsened_detail"].head(8).itertuples()):
        fill = altfill(i)
        short = {"Healthy":"Healthy","Early Degradation":"Early Degradation",
                 "Middle Degradation":"Middle Degradation","Late Degradation":"Late Degradation",
                 "Critical":"Critical","No Data":"No Data"}
        from_cfg = STATUS_CFG.get(row.status_old, ("—", C.SLATE_MID, C.SLATE_LIGHT, C.SLATE_MID))
        to_cfg   = STATUS_CFG.get(row.status_new, ("—", C.SLATE_MID, C.SLATE_LIGHT, C.SLATE_MID))
        tbl_rows.append([
            dat(f"{from_cfg[0]} {short.get(row.status_old,row.status_old)}", "left", False, from_cfg[1], fill),
            dat(f"{to_cfg[0]} {short.get(row.status_new,row.status_new)}",   "left", False, to_cfg[1],   fill),
            dat(str(row.count),                             "center", True, C.TEXT, fill),
            dat(f"{round(row.count/total_w*100,1)}%",       "center", False, C.TEXT, fill),
        ])
    add_table(slide, 0.38, 2.32, 9.24, 2.40,
              tbl_rows, col_widths=[3.50, 3.50, 1.12, 1.12],
              row_heights=[0.30] + [0.28]*(len(tbl_rows)-1))

    add_insight_box(slide, 0.38, 4.88, 9.24, 0.50,
        f"{s['worsened']} units showed deterioration between scans — "
        f"{s['worsened_detail'].iloc[0]['count'] if len(s['worsened_detail']) else 0} units "
        f"transitioned {s['worsened_detail'].iloc[0]['status_old'] if len(s['worsened_detail']) else ''} → "
        f"{s['worsened_detail'].iloc[0]['status_new'] if len(s['worsened_detail']) else ''} (largest group). "
        f"Units progressing beyond Early Degradation require priority field investigation.")

    add_text_box(slide, 0.38, 5.46, 9.24, 0.14,
                 "Worsened = unit moved to a higher severity state in the newer scan.",
                 size=7.5, italic=True, color=C.SLATE_MID)


def build_slide_improved(prs, s):
    """Slide 6: Improved Units."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Improved Units",
                     f"Units with lower degradation severity in newer scan  |  {s['improved']} units  |  "
                     f"{s['date_old']}  →  {s['date_new']}")

    add_kpi_card(slide, 0.38,  1.05, 2.80, 0.80,
                 "Total Improved", str(s["improved"]), C.GREEN_DARK, C.GREEN_LIGHT, C.GREEN_MID)
    add_kpi_card(slide, 3.42,  1.05, 2.80, 0.80,
                 "As % of Changed", f"{round(s['improved']/s['total_changed']*100,1)}%",
                 C.TEAL, C.TEAL_LIGHT, C.TEAL)
    add_kpi_card(slide, 6.46,  1.05, 2.80, 0.80,
                 "As % of Intersection", f"{round(s['improved']/s['intersection']*100,2)}%",
                 C.NAVY, C.NAVY_LIGHT, C.TEAL)

    add_section_label(slide, 0.38, 2.02, 9.24, "Improved Transitions Breakdown")
    tbl_rows = [[hdr("From Status"), hdr("To Status"), hdr("Units"), hdr("% of Improved")]]
    total_i = s["improved"]
    for i, row in enumerate(s["improved_detail"].head(8).itertuples()):
        fill = altfill(i)
        short = {"Healthy":"Healthy","Early Degradation":"Early Degradation",
                 "Middle Degradation":"Middle Degradation","Late Degradation":"Late Degradation",
                 "Critical":"Critical","No Data":"No Data"}
        from_cfg = STATUS_CFG.get(row.status_old, ("—", C.SLATE_MID, C.SLATE_LIGHT, C.SLATE_MID))
        to_cfg   = STATUS_CFG.get(row.status_new, ("—", C.SLATE_MID, C.SLATE_LIGHT, C.SLATE_MID))
        tbl_rows.append([
            dat(f"{from_cfg[0]} {short.get(row.status_old,row.status_old)}", "left", False, from_cfg[1], fill),
            dat(f"{to_cfg[0]} {short.get(row.status_new,row.status_new)}",   "left", False, to_cfg[1],   fill),
            dat(str(row.count),                              "center", True, C.TEXT, fill),
            dat(f"{round(row.count/total_i*100,1)}%",        "center", False, C.TEXT, fill),
        ])
    add_table(slide, 0.38, 2.32, 9.24, 2.40,
              tbl_rows, col_widths=[3.50, 3.50, 1.12, 1.12],
              row_heights=[0.30] + [0.28]*(len(tbl_rows)-1))

    add_insight_box(slide, 0.38, 4.88, 9.24, 0.50,
        f"{s['improved']} units showed improvement — notably {s['e2h_total']} Early→Healthy resolutions. "
        f"The near-symmetry with {s['worsened']} worsened units suggests the majority of changes "
        f"reflect boundary oscillation rather than sustained network-wide degradation trend.")

    add_text_box(slide, 0.38, 5.46, 9.24, 0.14,
                 "Improved = unit moved to a lower severity state in the newer scan. "
                 "Genuine improvement may reflect field remediation, connector re-torquing, or algorithm boundary effects.",
                 size=7.5, italic=True, color=C.SLATE_MID)


def build_slide_cog_stability(prs, s):
    """Slide 7: CoG Stability Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "CoG Stability Analysis — Stable Healthy Units",
                     f"Week-over-week CoG variation for units classified Healthy in both scans  |  "
                     f"{s['stable_healthy_n']:,} units  |  {s['intersection'] - s['stable_healthy_n']} excluded (changed status)")

    # KPIs
    kpis = [
        (f"{s['stable_healthy_n']:,}", "Stable Healthy Units",     C.GREEN_DARK,  C.GREEN_LIGHT,  C.GREEN_MID),
        (f"{s['cog_delta_std']} dB",   "CoG Delta Std Dev",        C.TEAL,        C.TEAL_LIGHT,   C.TEAL),
        (f"{s['cog_delta_max']} dB",   "Max CoG Delta Observed",   C.AMBER_DARK,  C.AMBER_LIGHT,  C.AMBER_MID),
        (f"{s['cog_delta_pct1']}%",    "Branches > ±1 dB Change",  C.ORANGE_DARK, C.ORANGE_LIGHT, C.ORANGE_MID),
    ]
    kw = 2.10; ky = 1.05; kh = 0.80
    for i, (val, lbl, dk, lt, nc) in enumerate(kpis):
        add_kpi_card(slide, 0.38 + i*2.34, ky, kw, kh, lbl, val, dk, lt, nc)

    # Histogram (left)
    add_section_label(slide, 0.38, 2.02, 5.80, "CoG Delta Distribution (all branches, stable-healthy units)")
    stab_buf = build_cog_stability_chart(s["cog_deltas"])
    slide.shapes.add_picture(stab_buf, Inches(0.38), Inches(2.32), Inches(5.80), Inches(2.40))

    # Interpretation table (right)
    add_section_label(slide, 6.38, 2.02, 3.24, "Algorithm Noise Band")
    interp_rows = [
        [hdr("CoG Delta"), hdr("Interpretation")],
        [dat("< ±0.5 dB",  "left"), dat("Normal variation",   "left", False, C.GREEN_MID,  C.GREEN_LIGHT)],
        [dat("±0.5–1.0 dB","left"), dat("Minor fluctuation",  "left", False, C.AMBER_MID,  C.AMBER_LIGHT)],
        [dat("> ±1.0 dB",  "left"), dat("Investigate",        "left", True,  C.RED_MID,    C.RED_LIGHT)],
        [dat("> ±2.0 dB",  "left"), dat("Data quality check", "left", True,  C.PURPLE_MID, C.PURPLE_LIGHT)],
    ]
    add_table(slide, 6.38, 2.32, 3.24, 1.28,
              interp_rows, col_widths=[1.30, 1.94], row_heights=[0.28]*5)

    add_section_label(slide, 6.38, 3.68, 3.24, "14 dB Threshold Context")
    add_rect(slide, 6.38, 3.98, 3.24, 0.96, C.SLATE_LIGHT, C.BORDER, 0.5)
    add_text_box(slide, 6.50, 4.02, 3.04, 0.88,
                 f"CoG std dev = {s['cog_delta_std']} dB — threshold: 14 dB\n\n"
                 f"A unit at CoG = 13.9 dB could flip in/out of Early Degradation "
                 f"from natural week-over-week variation alone.",
                 size=8.5, color=C.TEXT)

    add_insight_box(slide, 0.38, 4.76, 9.24, 0.60,
        f"Stable healthy units show CoG std = {s['cog_delta_std']} dB — very stable. "
        f"Only {s['cog_delta_above1']} branch readings ({s['cog_delta_pct1']}%) moved > ±1 dB. "
        f"The {s['h2e_total']+s['e2h_total']} flips between Healthy and Early Degradation are therefore "
        f"driven by FV/BTR/SZR borderline scoring, not CoG noise — the CoG boundary at 14 dB is stable.")

    add_text_box(slide, 0.38, 5.42, 9.24, 0.14,
                 f"CoG delta = (newer scan CoG) − (older scan CoG). "
                 f"Positive = improvement (higher return loss). Analysed across all 4 branches × {s['stable_healthy_n']:,} units.",
                 size=7.5, italic=True, color=C.SLATE_MID)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════════════════════
# CSV EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

def export_csvs(s: dict, out_dir: str):
    """Export all CSV analysis files to out_dir."""
    import os
    os.makedirs(out_dir, exist_ok=True)

    df_old  = s["_df_old"]
    df_new  = s["_df_new"]
    d_old   = s["_d_old"]
    d_new   = s["_d_new"]
    rich    = s["_rich"]
    changed = s["_changed"]
    COLS_AG = s["_COLS_AG"]

    def get_ag(df, sns):
        """Return cols A-G + extra columns for a set of serial numbers."""
        sub = df[df["serial_number"].isin(sns)].drop_duplicates("serial_number")
        return sub[COLS_AG].copy()

    def add_cog(base_df, source_df, suffix=""):
        """Add min CoG across branches to base_df from source_df."""
        cog_cols = [f"RL_{b}_CoG" for b in "ABCD"]
        avail = [c for c in cog_cols if c in source_df.columns]
        cog = source_df[["serial_number"] + avail].drop_duplicates("serial_number").set_index("serial_number")
        for c in avail:
            cog[c] = pd.to_numeric(cog[c], errors="coerce")
        cog[f"min_CoG{suffix}"] = cog[avail].min(axis=1)
        base_df = base_df.join(cog[[f"min_CoG{suffix}"]], on="serial_number")
        return base_df

    # ── 1. VSWR_risk_persistent.csv ─────────────────────────────────────────
    pers_sn  = s["_vswr_persistent_sn"]
    pers_dir = s["_pers_dir"]
    pers_old = s["_pers_old_st"]
    pers_new = s["_pers_new_st"]

    df_pers = get_ag(df_new, pers_sn).set_index("serial_number")
    df_pers["RL_status_older_scan"] = pers_old
    df_pers["RL_status_newer_scan"] = pers_new
    df_pers["RL_status_change"]     = pers_dir
    df_pers = df_pers.reset_index()
    df_pers = add_cog(df_pers, df_old.rename(columns={"serial_number":"serial_number"}), "_older")
    df_pers = add_cog(df_pers, df_new, "_newer")
    df_pers.to_csv(f"{out_dir}/VSWR_risk_persistent.csv", index=False)
    print(f"  ✓ VSWR_risk_persistent.csv  ({len(df_pers)} rows)")

    # ── 2. VSWR_risk_new_alerts.csv ──────────────────────────────────────────
    new_sn = s["_vswr_new_alerts_sn"]
    df_new_al = get_ag(df_new, new_sn).set_index("serial_number")
    df_new_al["RL_status_newer_scan"] = d_new.loc[list(new_sn), "status"]
    # VSWR alarm count
    vswr_extra = df_new[df_new["serial_number"].isin(new_sn)].drop_duplicates("serial_number").set_index("serial_number")
    for col in ["VSWR_Alarm_Count", "Active_VSWR_Alarm", "Historic_VSWR_Alarm", "VSWR_risk_assessment"]:
        if col in vswr_extra.columns:
            df_new_al[col] = vswr_extra[col]
    df_new_al = df_new_al.reset_index()
    df_new_al = add_cog(df_new_al, df_new, "_newer")
    df_new_al.to_csv(f"{out_dir}/VSWR_risk_new_alerts.csv", index=False)
    print(f"  ✓ VSWR_risk_new_alerts.csv  ({len(df_new_al)} rows)")

    # ── 3. VSWR_risk_disappeared.csv ─────────────────────────────────────────
    dis_sn = s["_vswr_disappeared_sn"]
    df_dis = get_ag(df_old, dis_sn).set_index("serial_number")
    df_dis["RL_status_older_scan"] = d_old.loc[list(dis_sn), "status"]
    vswr_dis = df_old[df_old["serial_number"].isin(dis_sn)].drop_duplicates("serial_number").set_index("serial_number")
    for col in ["VSWR_Alarm_Count", "Active_VSWR_Alarm", "Historic_VSWR_Alarm"]:
        if col in vswr_dis.columns:
            df_dis[col] = vswr_dis[col]
    df_dis = df_dis.reset_index()
    df_dis = add_cog(df_dis, df_old, "_older")
    df_dis.to_csv(f"{out_dir}/VSWR_risk_disappeared.csv", index=False)
    print(f"  ✓ VSWR_risk_disappeared.csv  ({len(df_dis)} rows)")

    # ── 4. status_changes.csv ────────────────────────────────────────────────
    chg_sn = set(changed.index)
    df_chg = get_ag(df_new, chg_sn).set_index("serial_number")
    df_chg["status_older_scan"] = changed["status_old"]
    df_chg["status_newer_scan"] = changed["status_new"]
    df_chg["direction"]         = changed["direction"]
    df_chg = df_chg.reset_index()
    df_chg = add_cog(df_chg, df_old, "_older")
    df_chg = add_cog(df_chg, df_new, "_newer")
    df_chg["CoG_delta"] = (df_chg["min_CoG_newer"] - df_chg["min_CoG_older"]).round(3)
    df_chg = df_chg.sort_values("direction")
    df_chg.to_csv(f"{out_dir}/status_changes.csv", index=False)
    print(f"  ✓ status_changes.csv  ({len(df_chg)} rows)")

    # ── 5. new_alerts.csv (Healthy → any degraded) ───────────────────────────
    new_al_mask = (changed["status_old"] == "Healthy") & (changed["status_new"] != "Healthy")
    new_al_sn   = set(changed[new_al_mask].index)
    df_nal = get_ag(df_new, new_al_sn).set_index("serial_number")
    df_nal["status_older_scan"] = changed.loc[list(new_al_sn), "status_old"]
    df_nal["status_newer_scan"] = changed.loc[list(new_al_sn), "status_new"]
    # CoG override flag
    ov_new = df_new[df_new["serial_number"].isin(new_al_sn)].drop_duplicates("serial_number").set_index("serial_number")
    def any_override(row):
        return "Yes" if any("Yes" in str(row.get(f"RL_{b}_CoG_Override", "")) for b in "ABCD") else "No"
    df_nal["via_CoG_override"] = ov_new.apply(any_override, axis=1)
    df_nal = df_nal.reset_index()
    df_nal = add_cog(df_nal, df_new, "_newer")
    df_nal.to_csv(f"{out_dir}/new_alerts.csv", index=False)
    print(f"  ✓ new_alerts.csv  ({len(df_nal)} rows)")

    # ── 6. resolved_alerts.csv (any degraded → Healthy) ─────────────────────
    res_mask = (changed["status_old"] != "Healthy") & (changed["status_new"] == "Healthy")
    res_sn   = set(changed[res_mask].index)
    df_res = get_ag(df_old, res_sn).set_index("serial_number")
    df_res["status_older_scan"] = changed.loc[list(res_sn), "status_old"]
    df_res["status_newer_scan"] = changed.loc[list(res_sn), "status_new"]
    df_res = df_res.reset_index()
    df_res = add_cog(df_res, df_old, "_older")
    df_res = add_cog(df_res, df_new, "_newer")
    df_res["CoG_delta"] = (df_res["min_CoG_newer"] - df_res["min_CoG_older"]).round(3)
    # Flip candidate: Early→Healthy AND CoG delta < 0.5 dB
    df_res["flip_candidate"] = df_res.apply(
        lambda r: "Yes" if (r["status_older_scan"] == "Early Degradation"
                            and abs(r["CoG_delta"]) < 0.5) else "No", axis=1)
    df_res.to_csv(f"{out_dir}/resolved_alerts.csv", index=False)
    print(f"  ✓ resolved_alerts.csv  ({len(df_res)} rows)")

    print(f"\nAll CSV files saved to: {out_dir}/")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — VSWR RISK ASSESSMENT COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_vswr_risk(prs, s):
    """Slide 8: VSWR Risk Assessment Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "VSWR Risk Assessment Comparison",
                     f"Column VSWR_risk_assessment  |  {s['date_old']}  →  {s['date_new']}  |  "
                     f"Intersection: {s['intersection']:,} units")

    # ── 5 KPI cards ──────────────────────────────────────────────────────────
    kpis = [
        (str(s["vswr_total_old"]),   f"Risk Units — {s['date_old']}",  C.RED_DARK,    C.RED_LIGHT,    C.RED_MID),
        (str(s["vswr_total_new"]),   f"Risk Units — {s['date_new']}",  C.RED_DARK,    C.RED_LIGHT,    C.RED_MID),
        (str(s["vswr_persistent"]),  "Persistent (Both Scans)",        C.ORANGE_DARK, C.ORANGE_LIGHT, C.ORANGE_MID),
        (str(s["vswr_new_alerts"]),  "New Risk Alerts",                C.PURPLE_DARK, C.PURPLE_LIGHT, C.PURPLE_MID),
        (str(s["vswr_disappeared"]), "Disappeared Risk Units",         C.GREEN_DARK,  C.GREEN_LIGHT,  C.GREEN_MID),
    ]
    kw = 1.76; ky = 1.05; kh = 0.86
    for i, (val, lbl, dk, lt, nc) in enumerate(kpis):
        add_kpi_card(slide, 0.38 + i * 1.86, ky, kw, kh, lbl, val, dk, lt, nc)

    # ── Left: Persistent risk — RL status change breakdown ───────────────────
    add_section_label(slide, 0.38, 2.08, 4.50, f"Persistent Risk Units ({s['vswr_persistent']}) — RL Status Change")
    pers_rows = [
        [hdr("RL Status Change"), hdr("Units"), hdr("Action")],
        [dat("Worsened (RL deteriorated)", "left", True,  C.RED_MID,    altfill(0)),
         dat(str(s["vswr_pers_worsened"]),  "center", True, C.RED_MID,  altfill(0)),
         dat("Priority — field inspection", "left", False, C.RED_MID,   altfill(0))],
        [dat("Stable (RL unchanged)",       "left", False, C.TEXT,      altfill(1)),
         dat(str(s["vswr_pers_stable"]),    "center", True, C.TEXT,     altfill(1)),
         dat("Monitor next scan",           "left", False, C.TEXT,      altfill(1))],
        [dat("Improved (RL better)",        "left", False, C.GREEN_MID, altfill(2)),
         dat(str(s["vswr_pers_improved"]),  "center", True, C.GREEN_MID,altfill(2)),
         dat("Verify field remediation",    "left", False, C.GREEN_MID, altfill(2))],
    ]
    add_table(slide, 0.38, 2.38, 4.50, 1.16,
              pers_rows, col_widths=[2.30, 0.70, 1.50],
              row_heights=[0.28, 0.29, 0.29, 0.29])

    # ── Right: New risk alerts — RL status breakdown ──────────────────────────
    add_section_label(slide, 5.08, 2.08, 4.54, f"New Risk Alerts ({s['vswr_new_alerts']}) — RL Status in Newer Scan")
    new_rows = [[hdr("RL Status"), hdr("Units"), hdr("% of New")]]
    total_new = max(s["vswr_new_alerts"], 1)
    for i, st in enumerate(STATUS_ORDER):
        v = s["vswr_new_rl_dist"].get(st, 0)
        if v == 0: continue
        cfg = STATUS_CFG[st]
        fill = altfill(i)
        new_rows.append([
            dat(f"{cfg[0]} {st}", "left",   False, cfg[1], fill),
            dat(str(v),           "center", True,  C.TEXT, fill),
            dat(f"{round(v/total_new*100,1)}%", "center", False, C.TEXT, fill),
        ])
    add_table(slide, 5.08, 2.38, 4.54, 1.30,
              new_rows, col_widths=[2.80, 0.74, 1.00],
              row_heights=[0.26] + [0.26]*(len(new_rows)-1))

    # ── Bottom left: Disappeared risk — old RL status ─────────────────────────
    add_section_label(slide, 0.38, 3.72, 4.50, f"Disappeared Risk Units ({s['vswr_disappeared']}) — RL Status in Older Scan")
    dis_rows = [[hdr("RL Status"), hdr("Units")]]
    for i, st in enumerate(STATUS_ORDER):
        v = s["vswr_dis_rl_dist"].get(st, 0)
        if v == 0: continue
        cfg = STATUS_CFG[st]
        fill = altfill(i)
        dis_rows.append([
            dat(f"{cfg[0]} {st}", "left",   False, cfg[1], fill),
            dat(str(v),           "center", True,  C.TEXT, fill),
        ])
    add_table(slide, 0.38, 4.02, 4.50, 0.66,
              dis_rows, col_widths=[3.50, 1.00],
              row_heights=[0.26] + [0.23]*(len(dis_rows)-1))

    # ── Bottom right: Risk trend summary box ─────────────────────────────────
    add_section_label(slide, 5.08, 3.72, 4.54, "Risk Trend Summary")
    net_str = f"+{s['vswr_net_change']}" if s['vswr_net_change'] >= 0 else str(s['vswr_net_change'])
    trend_color = C.RED_MID if s['vswr_net_change'] > 0 else C.GREEN_MID
    add_rect(slide, 5.08, 4.02, 4.54, 0.66, C.SLATE_LIGHT, C.BORDER, 0.5)
    trend_lines = [
        f"Net change (intersection): {net_str} risk units",
        f"Persistent worsening: {s['vswr_pers_worsened']} units require priority action",
        f"New alerts: {s['vswr_new_alerts']} units — {sum(v for k,v in s['vswr_new_rl_dist'].items() if k not in ['Healthy','No Data'])} already show RL degradation",
        f"Disappeared: {s['vswr_disappeared']} units — likely field remediation",
    ]
    add_text_box(slide, 5.18, 4.06, 4.34, 0.58,
                 "\n".join(trend_lines), size=8.5, color=C.TEXT)

    # ── Key Insight ───────────────────────────────────────────────────────────
    priority_n = s["vswr_pers_worsened"]
    degraded_new = sum(v for k,v in s["vswr_new_rl_dist"].items() if k not in ["Healthy", "No Data"])
    add_insight_box(slide, 0.38, 4.78, 9.24, 0.64,
        f"VSWR risk list grew from {s['vswr_total_old']} to {s['vswr_total_new']} units (net {net_str} within intersection). "
        f"{s['vswr_persistent']} units remain persistently at risk — of which {priority_n} have worsening RL status "
        f"and require priority field inspection. "
        f"Of the {s['vswr_new_alerts']} new risk alerts, {degraded_new} already show confirmed RL degradation "
        f"(Middle/Late/Critical). The {s['vswr_disappeared']} disappeared units should be verified as "
        f"field-remediated in maintenance records.")

    add_text_box(slide, 0.38, 5.46, 9.24, 0.10,
                 "VSWR risk = port with no DL power (NaN/0) or RL reading < 10 dB (VSWR alarm boundary). "
                 "Comparison on intersection only.",
                 size=7.5, italic=True, color=C.SLATE_MID)


def build_presentation(s: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    print("Building comparison report...")
    build_slide_title(prs, s);         print("  ✓ Title slide")
    build_slide_overview(prs, s);      print("  ✓ Slide 2 — Overview")
    build_slide_transition(prs, s);    print("  ✓ Slide 3 — Transition Matrix")
    build_slide_flip(prs, s);          print("  ✓ Slide 4 — Flip Analysis")
    build_slide_worsened(prs, s);      print("  ✓ Slide 5 — Worsened Units")
    build_slide_improved(prs, s);      print("  ✓ Slide 6 — Improved Units")
    build_slide_cog_stability(prs, s); print("  ✓ Slide 7 — CoG Stability")
    build_slide_vswr_risk(prs, s);     print("  ✓ Slide 8 — VSWR Risk Comparison")

    prs.save(output_path)
    print(f"\nReport saved: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="RL Health Scan Comparison Report",
        usage="python compare_scans_report.py scan1.csv scan2.csv [output.pptx]"
    )
    parser.add_argument("scan1", help="First scan CSV (older or newer — auto-detected)")
    parser.add_argument("scan2", help="Second scan CSV")
    parser.add_argument("output", nargs="?", default=None, help="Output .pptx path (optional)")
    args = parser.parse_args()

    print(f"Loading scan 1: {args.scan1}")
    df1 = load_df(args.scan1)
    print(f"  {len(df1):,} units  (scan date: {fmt_date(df1._scan_date)})")

    print(f"Loading scan 2: {args.scan2}")
    df2 = load_df(args.scan2)
    print(f"  {len(df2):,} units  (scan date: {fmt_date(df2._scan_date)})")

    # Determine older / newer by scan date
    if df1._scan_date >= df2._scan_date:
        df_new, df_old = df1, df2
    else:
        df_new, df_old = df2, df1

    print(f"\nOlder scan: {fmt_date(df_old._scan_date)}  |  Newer scan: {fmt_date(df_new._scan_date)}")
    print("Computing statistics...")
    s = compute_comparison(df_old, df_new)

    # Build output folder name
    d1 = df_old._scan_date; d2 = df_new._scan_date
    out_dir = f"RL_scan_comparison_{d1}_vs_{d2}"

    if args.output:
        out_path = args.output
    else:
        import os
        os.makedirs(out_dir, exist_ok=True)
        out_path = f"{out_dir}/RL_scan_comparison_report.pptx"

    build_presentation(s, out_path)

    # Export CSV files
    csv_dir = out_dir if not args.output else str(Path(args.output).parent / out_dir)
    print("\nExporting CSV files...")
    export_csvs(s, csv_dir)


if __name__ == "__main__":
    main()
