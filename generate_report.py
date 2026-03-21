"""
generate_report.py
==================
Auto-generates the Radio 4432 B28 Connector RL Health Analysis Report (6 slides)
from the monthly network scan CSV export.

Usage:
    python3 generate_report.py --input RRU_BXP_RL50_Health_Result.csv
    python3 generate_report.py --input data.csv --output MyReport.pptx

Dependencies:
    pip install python-pptx matplotlib
"""

import argparse
import io
import sys
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ═══════════════════════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════════

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Colour palette (RGB tuples)
class C:
    NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
    NAVY_MID    = RGBColor(0x2C, 0x4F, 0x7A)
    NAVY_LIGHT  = RGBColor(0xEB, 0xF0, 0xF7)
    TEAL        = RGBColor(0x08, 0x91, 0xB2)
    TEAL_DARK   = RGBColor(0x0E, 0x74, 0x90)
    TEAL_LIGHT  = RGBColor(0xE0, 0xF2, 0xFE)
    GREEN_DARK  = RGBColor(0x2D, 0x6A, 0x2D)
    GREEN_MID   = RGBColor(0x16, 0xA3, 0x4A)
    GREEN_LIGHT = RGBColor(0xEB, 0xF5, 0xEB)
    AMBER_DARK  = RGBColor(0x92, 0x40, 0x0E)
    AMBER_MID   = RGBColor(0xD9, 0x77, 0x06)
    AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
    ORANGE_DARK = RGBColor(0x7C, 0x2D, 0x12)
    ORANGE_MID  = RGBColor(0xEA, 0x58, 0x0C)
    ORANGE_LIGHT= RGBColor(0xFF, 0xF0, 0xE6)
    RED_DARK    = RGBColor(0x7F, 0x1D, 0x1D)
    RED_MID     = RGBColor(0xDC, 0x26, 0x26)
    RED_LIGHT   = RGBColor(0xFE, 0xF2, 0xF2)
    PURPLE_DARK = RGBColor(0x4C, 0x1D, 0x95)
    PURPLE_MID  = RGBColor(0x7C, 0x3A, 0xED)
    PURPLE_LIGHT= RGBColor(0xF3, 0xE8, 0xFF)
    SLATE_DARK  = RGBColor(0x37, 0x41, 0x51)
    SLATE_MID   = RGBColor(0x64, 0x74, 0x8B)
    SLATE_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
    BORDER      = RGBColor(0xE2, 0xE8, 0xF0)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT        = RGBColor(0x1E, 0x29, 0x3B)
    BG          = RGBColor(0xF4, 0xF6, 0xFA)

# Status config: (label, emoji, dark, light, num_color)
STATUS_CFG = {
    "Healthy":            ("Healthy",            "✅", C.GREEN_DARK,  C.GREEN_LIGHT,  C.GREEN_MID),
    "Early Degradation":  ("Early Degradation",  "⚠️", C.AMBER_DARK,  C.AMBER_LIGHT,  C.AMBER_MID),
    "Middle Degradation": ("Middle Degradation", "🔶", C.ORANGE_DARK, C.ORANGE_LIGHT, C.ORANGE_MID),
    "Late Degradation":   ("Late Degradation",   "🔴", C.RED_DARK,    C.RED_LIGHT,    C.RED_MID),
    "Critical":           ("Critical / Failed",  "🚨", C.PURPLE_DARK, C.PURPLE_LIGHT, C.PURPLE_MID),
    "No Data":            ("No Data",            "—",  C.SLATE_DARK,  C.SLATE_LIGHT,  C.SLATE_MID),
}

# ═══════════════════════════════════════════════════════════════════════════════
# DATA LOADING & COMPUTATION
# ═══════════════════════════════════════════════════════════════════════════════

def load_data(csv_path: str) -> pd.DataFrame:
    df = pd.read_csv(csv_path)
    # Parse VSWR columns
    for c in ["VSWR_1_dB", "VSWR_2_dB", "VSWR_3_dB", "VSWR_4_dB"]:
        df[c] = pd.to_numeric(df[c], errors="coerce")
    # Parse deployment date
    df["deploy_date"] = pd.to_datetime(
        df["Unit_Field_Deployment_Date"], format="%y-%m-%d", errors="coerce"
    )
    # Extract clean status
    df["status"] = df["Unit_Worst_Status"].str.extract(
        r"(Healthy|Early Degradation|Middle Degradation|Late Degradation|Critical|No Data)",
        expand=False,
    )
    return df


def compute_stats(df: pd.DataFrame) -> dict:
    s = {}
    total = len(df)
    s["total_units"] = total

    # Scan date: derive from latest RL timestamp in data
    scan_date_str = "Feb 2026"
    for _col in ["RL_A_time", "RL_B_time", "RL_C_time", "RL_D_time"]:
        if _col in df.columns:
            _sample = df[_col].dropna().astype(str)
            _sample = _sample[_sample.str.match(r"^\d{6}")]
            if len(_sample) > 0:
                try:
                    _dt = datetime.strptime(_sample.max()[:6], "%y%m%d")
                    scan_date_str = _dt.strftime("%b %Y")
                    break
                except Exception:
                    pass
    s["scan_date"] = scan_date_str

    # Count unique source files (sites)
    s["total_sites"] = df["source_file"].nunique() if "source_file" in df.columns else "N/A"

    # ── PAGE 1: status counts ─────────────────────────────────────────────
    counts = df["status"].value_counts()
    s["status_counts"] = {k: int(counts.get(k, 0)) for k in STATUS_CFG}
    healthy_n = s["status_counts"]["Healthy"]
    s["healthy_pct"] = round(healthy_n / total * 100, 1)
    alert_keys = ["Early Degradation", "Middle Degradation", "Late Degradation", "Critical"]
    s["alert_total"] = sum(s["status_counts"][k] for k in alert_keys)

    # ── PAGE 2: branches affected per unit ───────────────────────────────
    alert_mask = df["status"].isin(alert_keys)
    alert_df = df[alert_mask].copy()

    def count_alert_branches(row):
        cnt = 0
        for br in ["A", "B", "C", "D"]:
            col = f"RL_{br}_Health"
            if col in row.index:
                h = str(row[col])
                if any(x in h for x in ["Early", "Middle", "Late", "Critical"]):
                    cnt += 1
        return cnt

    if len(alert_df) > 0:
        alert_df = alert_df.copy()
        alert_df["branch_count"] = alert_df.apply(count_alert_branches, axis=1)
        alert_df["sc"] = alert_df["status"]
        branch_cross = (
            alert_df.groupby(["sc", "branch_count"])
            .size()
            .unstack(fill_value=0)
            .reindex(columns=[1, 2, 3, 4], fill_value=0)
        )
        branch_cross["Total"] = branch_cross.sum(axis=1)
        s["branch_table"] = branch_cross
        s["branch_totals"] = {
            1: int(branch_cross[1].sum()),
            2: int(branch_cross[2].sum()),
            3: int(branch_cross[3].sum()),
            4: int(branch_cross[4].sum()),
            "Total": int(branch_cross["Total"].sum()),
        }
        one_branch = int(branch_cross[1].sum())
        grand_total = int(branch_cross["Total"].sum())
        s["one_branch_pct"] = round(one_branch / grand_total * 100) if grand_total > 0 else 0
    else:
        s["branch_table"] = pd.DataFrame()
        s["branch_totals"] = {1: 0, 2: 0, 3: 0, 4: 0, "Total": 0}
        s["one_branch_pct"] = 0

    # ── PAGE 3: active VSWR alarm ────────────────────────────────────────
    if "Active_VSWR_Alarm" in df.columns:
        active_mask = df["Active_VSWR_Alarm"].notna() & ~df["Active_VSWR_Alarm"].str.contains(
            "No active", case=False, na=True
        )
    else:
        active_mask = pd.Series([False] * total, index=df.index)

    s["active_alarm_total"] = int(active_mask.sum())
    s["no_alarm_total"] = total - s["active_alarm_total"]
    s["alarm_rate_pct"] = round(s["active_alarm_total"] / total * 100, 1)

    alarm_df = df[active_mask].copy()
    alarm_by_status = alarm_df["status"].value_counts()
    s["alarm_by_status"] = {k: int(alarm_by_status.get(k, 0)) for k in STATUS_CFG}
    degraded_with_alarm = sum(
        s["alarm_by_status"].get(k, 0) for k in alert_keys
    )
    s["degraded_with_alarm"] = degraded_with_alarm

    # ── PAGE 4: VSWR risk ────────────────────────────────────────────────
    vswr_cols = ["VSWR_1_dB", "VSWR_2_dB", "VSWR_3_dB", "VSWR_4_dB"]
    vswr_all = pd.concat([df[c].dropna() for c in vswr_cols])
    vswr_valid = vswr_all[(vswr_all > 0) & (vswr_all <= 40)]
    s["vswr_valid_count"] = len(vswr_valid)
    s["vswr_mean"] = round(float(vswr_valid.mean()), 1) if len(vswr_valid) > 0 else 0
    s["vswr_readings"] = vswr_valid.values  # for histogram

    risk_counts = df["VSWR_risk_assessment"].value_counts() if "VSWR_risk_assessment" in df.columns else {}
    s["risk_units"]     = int(risk_counts.get("Risk unit", 0))
    s["non_risk_units"] = int(risk_counts.get("Non-risk unit", 0))
    s["risk_pct"]       = round(s["risk_units"] / total * 100, 1)
    s["non_risk_pct"]   = round(s["non_risk_units"] / total * 100, 1)

    # Sub-14 dB port analysis (CoG boundary)
    s["ports_below_14"] = int((vswr_valid < 14).sum())
    s["ports_below_14_pct"] = round(s["ports_below_14"] / len(vswr_valid) * 100, 1) if len(vswr_valid) > 0 else 0

    # How many of those sub-14 ports come from degraded vs healthy units
    alert_mask2 = df["status"].isin(alert_keys)
    alert_vswr = pd.concat([df.loc[alert_mask2, c].dropna() for c in vswr_cols])
    alert_vswr_valid = alert_vswr[(alert_vswr > 0) & (alert_vswr <= 40)]
    s["ports_below_14_degraded"] = int((alert_vswr_valid < 14).sum())
    s["ports_below_14_healthy"]  = s["ports_below_14"] - s["ports_below_14_degraded"]

    # Average degraded ports per degraded unit
    s["avg_degraded_ports"] = round(s["ports_below_14_degraded"] / s["alert_total"], 2) if s["alert_total"] > 0 else 0

    # ── PAGE 5: deployment age cohorts ───────────────────────────────────
    scan_ts = pd.Timestamp("2026-02-27")
    df["age_months"] = (scan_ts - df["deploy_date"]).dt.days / 30.44

    cohort_bins   = [0, 6, 12, 24, 9999]
    cohort_labels = ["< 6 months", "6 – 12 months", "1 – 2 years", "> 2 years"]
    df["cohort"] = pd.cut(df["age_months"], bins=cohort_bins, labels=cohort_labels).astype(str)
    df.loc[df["deploy_date"].isna(), "cohort"] = "Unknown"

    cohort_order = ["< 6 months", "6 – 12 months", "1 – 2 years", "> 2 years", "Unknown"]
    cohort_rows = []
    for coh in cohort_order:
        sub = df[df["cohort"] == coh]
        n = len(sub)
        sc = sub["status"].value_counts()
        healthy = int(sc.get("Healthy", 0))
        early   = int(sc.get("Early Degradation", 0))
        middle  = int(sc.get("Middle Degradation", 0))
        late    = int(sc.get("Late Degradation", 0))
        crit    = int(sc.get("Critical", 0))
        nodata  = int(sc.get("No Data", 0))
        alert   = early + middle + late + crit
        rate    = round(alert / n * 100, 1) if n > 0 else 0
        small   = n < 700
        cohort_rows.append({
            "cohort": coh, "n": n, "healthy": healthy,
            "early": early, "middle": middle, "late": late,
            "critical": crit, "nodata": nodata,
            "alert": alert, "rate": rate, "small": small,
        })
    s["cohort_rows"] = cohort_rows

    return s


# ═══════════════════════════════════════════════════════════════════════════════
# DRAWING HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def rgb(color: RGBColor):
    return color

def inches(*args):
    return [Inches(a) for a in args]

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width_pt=0):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt) if line_width_pt else Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size, bold=False, color=None,
                 align=PP_ALIGN.LEFT, italic=False, font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txb


def add_rich_text_box(slide, x, y, w, h, runs, font_name="Calibri", align=PP_ALIGN.LEFT, wrap=True):
    """runs = list of (text, font_size, bold, color, italic, breakline)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    for i, (text, fsize, bold, color, italic) in enumerate(runs):
        run = para.add_run()
        run.text = text
        run.font.size = Pt(fsize)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
        if color:
            run.font.color.rgb = color
    return txb


def add_slide_header(slide, title, subtitle=None):
    """Full-width navy header band with optional subtitle band."""
    # Main navy bar
    add_rect(slide, 0, 0, 10, 0.72, C.NAVY)
    add_text_box(slide, 0.38, 0.04, 9.24, 0.64, title,
                 font_size=20, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_rect(slide, 0, 0.72, 10, 0.30, C.NAVY_MID)
        add_text_box(slide, 0.38, 0.74, 9.24, 0.26, subtitle,
                     font_size=9.5, color=RGBColor(0xC8, 0xD8, 0xEC), align=PP_ALIGN.LEFT)


def add_section_label(slide, x, y, w, text):
    add_rect(slide, x, y, w, 0.26, C.TEAL_DARK)
    add_text_box(slide, x + 0.10, y + 0.02, w - 0.12, 0.22, text,
                 font_size=10, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)


def add_insight_box(slide, x, y, w, h, text):
    add_rect(slide, x, y, w, h, C.SLATE_LIGHT, C.BORDER)
    add_rect(slide, x, y, 0.055, h, C.NAVY)
    # "Key Insight" bold label + body text
    txb = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.06),
                                   Inches(w - 0.20), Inches(h - 0.10))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    # First paragraph: bold label
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = "Key Insight"
    r1.font.bold = True
    r1.font.size = Pt(10)
    r1.font.name = "Calibri"
    r1.font.color.rgb = C.NAVY
    # Second paragraph: body
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = text
    r2.font.bold = False
    r2.font.size = Pt(10)
    r2.font.name = "Calibri"
    r2.font.color.rgb = C.TEXT


def add_status_card(slide, x, y, w, h, label, count_str, pct_str, dark, light, num_color):
    """Status KPI card: dark header band + light body with large number + pct."""
    hdr_h = h * 0.30
    body_h = h - hdr_h
    # Header
    add_rect(slide, x, y, w, hdr_h, dark)
    add_text_box(slide, x, y + 0.01, w, hdr_h - 0.02, label,
                 font_size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    # Body
    add_rect(slide, x, y + hdr_h, w, body_h, light, C.BORDER, 0.5)
    # Count
    add_text_box(slide, x, y + hdr_h + 0.04, w, body_h * 0.55,
                 count_str, font_size=28, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    # Pct
    add_text_box(slide, x, y + hdr_h + body_h * 0.60, w, body_h * 0.36,
                 pct_str, font_size=11, bold=False, color=num_color, align=PP_ALIGN.CENTER)


def add_kpi_card(slide, x, y, w, h, label, value, dark, light, num_color):
    """Simple KPI card (for pages 3/4)."""
    hdr_h = h * 0.30
    body_h = h - hdr_h
    add_rect(slide, x, y, w, hdr_h, dark)
    add_text_box(slide, x, y + 0.01, w, hdr_h - 0.02, label,
                 font_size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, y + hdr_h, w, body_h, light, C.BORDER, 0.5)
    add_text_box(slide, x, y + hdr_h + 0.06, w, body_h * 0.70,
                 str(value), font_size=28, bold=True, color=num_color, align=PP_ALIGN.CENTER)


# ── TABLE BUILDER ────────────────────────────────────────────────────────────

def add_table(slide, x, y, w, h, rows, col_widths, row_heights=None):
    """
    rows: list of list of dicts with keys:
        text, font_size, bold, color, align, fill, italic
    col_widths: list of floats (inches), must sum to w
    row_heights: list of floats (inches), must match len(rows)
    """
    num_rows = len(rows)
    num_cols = len(rows[0])
    tbl = slide.shapes.add_table(num_rows, num_cols,
                                  Inches(x), Inches(y), Inches(w), Inches(h)).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    # Set row heights
    if row_heights:
        for ri, rh in enumerate(row_heights):
            tbl.rows[ri].height = Inches(rh)

    for ri, row in enumerate(rows):
        for ci, cell_cfg in enumerate(row):
            cell = tbl.cell(ri, ci)
            # Fill
            if cell_cfg.get("fill"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_cfg["fill"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C.WHITE

            # Remove cell borders individually (keep outer border clean)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            align_map = {
                "center": PP_ALIGN.CENTER,
                "left":   PP_ALIGN.LEFT,
                "right":  PP_ALIGN.RIGHT,
            }
            para.alignment = align_map.get(cell_cfg.get("align", "center"), PP_ALIGN.CENTER)

            run = para.add_run()
            run.text = str(cell_cfg.get("text", ""))
            run.font.size = Pt(cell_cfg.get("font_size", 10.5))
            run.font.bold = cell_cfg.get("bold", False)
            run.font.italic = cell_cfg.get("italic", False)
            run.font.name = "Calibri"
            color = cell_cfg.get("color", C.TEXT)
            if color:
                run.font.color.rgb = color

            # Vertical centering
            cell.vertical_anchor = 3  # middle

    return tbl


def hdr_cell(text, align="center"):
    return {"text": text, "font_size": 10.5, "bold": True,
            "color": C.WHITE, "fill": C.NAVY, "align": align}

def data_cell(text, align="center", bold=False, color=None, fill=None, italic=False, font_size=10.5):
    return {"text": text, "font_size": font_size, "bold": bold,
            "color": color or C.TEXT, "fill": fill or C.WHITE,
            "align": align, "italic": italic}

def row_bg(i):
    return C.WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)


# ── HISTOGRAM ────────────────────────────────────────────────────────────────

def build_histogram(vswr_values) -> bytes:
    fig, ax = plt.subplots(figsize=(6.0, 3.2), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    bins = np.arange(0, 41, 1)
    ax.hist(vswr_values, bins=bins, color="#0891B2", edgecolor="white", linewidth=0.4)

    ax.axvline(x=14, color="#DC2626", linewidth=1.8, linestyle="--", zorder=5)
    ylim = ax.get_ylim()
    ax.text(0.5, ylim[1] * 0.93, "RL degradation\nthreshold (14 dB)",
            color="#DC2626", fontsize=8, fontweight="bold", ha="left", va="top")
    ax.annotate("", xy=(14, ylim[1]*0.80), xytext=(9, ylim[1]*0.80),
                arrowprops=dict(arrowstyle="-|>", color="#DC2626", lw=1.2))

    ax.set_xlabel("VSWR (dB)", fontsize=9.5, color="#1E293B", labelpad=5)
    ax.set_ylabel("Port count", fontsize=9.5, color="#1E293B", labelpad=5)
    ax.set_title("VSWR Port Measurement Distribution", fontsize=10.5,
                 fontweight="bold", color="#1E293B", pad=8)
    ax.set_xlim(0, 40)
    ax.tick_params(colors="#64748B", labelsize=8.5)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#E2E8F0")
    ax.spines["bottom"].set_color("#E2E8F0")
    ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.6, zorder=0)
    ax.set_axisbelow(True)
    ax.set_xticks(range(0, 41, 5))

    plt.tight_layout(pad=0.6)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_title(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title and setattr(slide.shapes.title, "text", "")

    # Full dark navy background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C.NAVY

    # Antenna unicode icon
    add_text_box(slide, 0.50, 0.28, 3.0, 0.70, '("A")',
                 font_size=28, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)

    # Main title
    add_text_box(slide, 0.50, 0.95, 9.0, 1.60,
                 "Radio 4432 B28 Connector RL Health\nAnalysis Report",
                 font_size=36, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)

    # Teal rule
    add_rect(slide, 0.50, 2.62, 8.60, 0.04, C.TEAL)

    # Subtitle
    add_text_box(slide, 0.50, 2.74, 9.0, 0.38,
                 f"Network-Wide Automated Degradation Detection  |  RRU4432B28  |  {s['scan_date']}",
                 font_size=12, bold=False, color=RGBColor(0xC8, 0xD8, 0xEC), align=PP_ALIGN.LEFT)

    # 4 KPI boxes
    kpis = [
        (f"{s['total_units']:,}",  "Units Scanned"),
        (f"{s['alert_total']:,}",  "Alert Units"),
        (f"{s['healthy_pct']}%",   "Network Healthy"),
        (f"{s['total_sites']:,}",  "Sites Covered"),
    ]
    box_w, box_h = 2.10, 0.88
    gap = 0.16
    start_x = 0.50
    box_y = 3.54

    for i, (val, lbl) in enumerate(kpis):
        bx = start_x + i * (box_w + gap)
        add_rect(slide, bx, box_y, box_w, box_h, RGBColor(0x2C, 0x4F, 0x7A), C.TEAL, 0.75)
        add_text_box(slide, bx, box_y + 0.06, box_w, box_h * 0.55,
                     val, font_size=22, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, bx, box_y + box_h * 0.60, box_w, box_h * 0.36,
                     lbl, font_size=10, bold=False,
                     color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, 0.50, 5.30, 5.0, 0.24,
                 "Confidential — Internal Management Report",
                 font_size=9, bold=False, color=RGBColor(0x6B, 0x82, 0x9A), align=PP_ALIGN.LEFT)


def build_slide_1(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.BG

    add_slide_header(slide, "Network Scan Results — Executive Summary",
                     f"{s['scan_date']}  |  {s['total_units']:,} units scanned across {s['total_sites']:,} sites  |  Model: RRU4432B28")

    # Hero KPI — total units
    add_rect(slide, 0.35, 1.12, 2.15, 1.85, C.NAVY)
    add_text_box(slide, 0.35, 1.22, 2.15, 1.10,
                 f"{s['total_units']:,}", font_size=36, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.35, 2.30, 2.15, 0.60,
                 "Total Units\nScanned", font_size=11, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

    # 6 status cards
    card_defs = [
        ("Healthy",            s["status_counts"]["Healthy"],            f"{s['healthy_pct']}%"),
        ("Early Degradation",  s["status_counts"]["Early Degradation"],  f"{s['status_counts']['Early Degradation']/s['total_units']*100:.1f}%"),
        ("Middle Degradation", s["status_counts"]["Middle Degradation"], f"{s['status_counts']['Middle Degradation']/s['total_units']*100:.1f}%"),
        ("Late Degradation",   s["status_counts"]["Late Degradation"],   f"{s['status_counts']['Late Degradation']/s['total_units']*100:.1f}%"),
        ("Critical",           s["status_counts"]["Critical"],           f"{s['status_counts']['Critical']/s['total_units']*100:.1f}%"),
        ("No Data",            s["status_counts"]["No Data"],            f"{s['status_counts']['No Data']/s['total_units']*100:.1f}%"),
    ]

    card_w, card_h = 2.45, 0.88
    gap = 0.10
    start_x = 2.68
    rows_y = [1.12, 2.10]

    for i, (key, count, pct) in enumerate(card_defs):
        cfg = STATUS_CFG[key]
        label = cfg[0]  # display label
        dark, light, num_col = cfg[2], cfg[3], cfg[4]
        col = i % 3
        row = i // 3
        cx = start_x + col * (card_w + gap)
        cy = rows_y[row]
        add_status_card(slide, cx, cy, card_w, card_h,
                        label, f"{count:,}", pct, dark, light, num_col)

    # Bar chart via matplotlib embedded as image
    fig, ax = plt.subplots(figsize=(9.3, 2.1), dpi=130)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    labels_bar = ["Healthy", "Early Degr.", "Middle Degr.", "Late Degr.", "Critical/Failed", "No Data"]
    values_bar = [s["status_counts"][k] for k in ["Healthy","Early Degradation","Middle Degradation","Late Degradation","Critical","No Data"]]
    colors_bar = ["#16A34A","#D97706","#EA580C","#DC2626","#7C3AED","#64748B"]
    bars = ax.bar(labels_bar, values_bar, color=colors_bar, width=0.55, edgecolor="white")
    for bar, val in zip(bars, values_bar):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 30,
                f"{val:,}", ha="center", va="bottom", fontsize=9, color="#1E293B", fontweight="bold")
    ax.set_ylim(0, max(values_bar) * 1.18)
    ax.tick_params(colors="#64748B", labelsize=9)
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#E2E8F0"); ax.spines["bottom"].set_color("#E2E8F0")
    ax.yaxis.grid(True, color="#E2E8F0", linewidth=0.5, zorder=0); ax.set_axisbelow(True)
    ax.set_ylabel("Unit count", fontsize=8.5, color="#64748B")
    plt.tight_layout(pad=0.3)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="white")
    plt.close(fig); buf.seek(0)
    slide.shapes.add_picture(buf, Inches(0.35), Inches(3.12), Inches(9.30), Inches(2.25))

    add_text_box(slide, 0.35, 5.40, 9.30, 0.20,
                 "Note: Classification based on weekly Event 50 RL Statistic Counter snapshot. Reflects connector RF path health, not real-time alarm status.",
                 font_size=8, italic=True, color=C.SLATE_MID)


def build_slide_2(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Degradation Units — Branch-Level Breakdown",
                     "Number of branches (out of 4) affected per degraded unit")

    bt = s["branch_table"]
    order = ["Early Degradation", "Middle Degradation", "Late Degradation", "Critical"]
    emojis = {"Early Degradation":"⚠️ ", "Middle Degradation":"🔶 ", "Late Degradation":"🔴 ", "Critical":"🚨 "}
    colors = {"Early Degradation": C.AMBER_MID, "Middle Degradation": C.ORANGE_MID,
              "Late Degradation": C.RED_MID, "Critical": C.PURPLE_MID}

    def dash(n): return "—" if n == 0 else str(n)

    rows = [[hdr_cell("Status", "left"),
             hdr_cell("1/4 Branches"), hdr_cell("2/4 Branches"),
             hdr_cell("3/4 Branches"), hdr_cell("4/4 Branches"), hdr_cell("Total")]]

    for i, key in enumerate(order):
        if key not in bt.index:
            r = {c: 0 for c in [1,2,3,4,"Total"]}
        else:
            r = bt.loc[key]
        bg_c = row_bg(i)
        col = colors[key]
        lbl = emojis[key] + STATUS_CFG[key][0]
        rows.append([
            data_cell(lbl,         "left",   True,  col,    bg_c),
            data_cell(dash(int(r.get(1,0))), "center", False, C.TEXT, bg_c),
            data_cell(dash(int(r.get(2,0))), "center", False, C.TEXT, bg_c),
            data_cell(dash(int(r.get(3,0))), "center", False, C.TEXT, bg_c),
            data_cell(dash(int(r.get(4,0))), "center", False, C.TEXT, bg_c),
            data_cell(str(int(r.get("Total",0))), "center", True, col, bg_c),
        ])

    tot = s["branch_totals"]
    rows.append([
        data_cell("Grand Total (Alert Units)", "left", True, C.NAVY, C.NAVY_LIGHT),
        data_cell(str(tot[1]),       "center", True, C.NAVY, C.NAVY_LIGHT),
        data_cell(str(tot[2]),       "center", True, C.NAVY, C.NAVY_LIGHT),
        data_cell(str(tot[3]),       "center", True, C.NAVY, C.NAVY_LIGHT),
        data_cell(str(tot[4]),       "center", True, C.NAVY, C.NAVY_LIGHT),
        data_cell(str(tot["Total"]), "center", True, C.NAVY, C.NAVY_LIGHT),
    ])

    add_table(slide, 0.90, 1.18, 8.20, 2.80, rows,
              col_widths=[2.80, 1.10, 1.10, 1.10, 1.10, 1.00],
              row_heights=[0.42, 0.50, 0.50, 0.50, 0.50, 0.42])

    pct = s["one_branch_pct"]
    one = tot[1]; grand = tot["Total"]
    add_insight_box(slide, 0.90, 4.14, 8.20, 0.88,
        f"{pct}% of degraded units ({one} out of {grand}) have only 1 out of 4 branches affected — "
        f"indicating early-stage, single-branch connector aging. Multi-branch degradation (2–4 branches) "
        f"in the same unit signals a more advanced failure requiring urgent action.")

    add_text_box(slide, 0.90, 5.40, 8.20, 0.20,
                 "* Branch health assessed independently. Unit worst status reflects the most severe branch result.",
                 font_size=8, italic=True, color=C.SLATE_MID)


def build_slide_3(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Active VSWR Alarm Analysis",
                     f"Units with VSWR alarm active at time of network scan — {s['scan_date']}")

    # 3 KPI cards
    kpi_defs = [
        (str(s["active_alarm_total"]),        "Units with Active VSWR Alarm", C.RED_DARK,   C.RED_LIGHT,   C.RED_MID),
        (f"{s['no_alarm_total']:,}",          "Units — No Active Alarm",      C.GREEN_DARK, C.GREEN_LIGHT, C.GREEN_MID),
        (f"{s['alarm_rate_pct']}%",           "Active Alarm Rate",             C.AMBER_DARK, C.AMBER_LIGHT, C.AMBER_MID),
    ]
    kw = 2.90; ky = 1.10; kh = 0.88
    for i, (val, lbl, dk, lt, nc) in enumerate(kpi_defs):
        add_kpi_card(slide, 0.38 + i * 3.12, ky, kw, kh, lbl, val, dk, lt, nc)

    add_section_label(slide, 0.38, 2.16, 9.24, "Active Alarm Units — RL Health Status Breakdown")

    # Build alarm table rows
    alarm_order = ["Critical", "Late Degradation", "Middle Degradation", "Early Degradation", "Healthy"]
    implications = {
        "Critical":           "RF path failure confirmed — immediate replacement",
        "Late Degradation":   "Severe degradation + active alarm — urgent inspection",
        "Middle Degradation": "Progressive degradation — schedule field maintenance",
        "Early Degradation":  "Early wear with alarm — escalate monitoring",
        "Healthy":            "Transient alarm — verify connection and monitor",
    }
    emojis3 = {"Critical":"🚨 ","Late Degradation":"🔴 ","Middle Degradation":"🔶 ",
               "Early Degradation":"⚠️ ","Healthy":"✅ "}
    colors3 = {"Critical": C.PURPLE_MID, "Late Degradation": C.RED_MID,
               "Middle Degradation": C.ORANGE_MID, "Early Degradation": C.AMBER_MID,
               "Healthy": C.GREEN_MID}

    alarm_rows = [[hdr_cell("RL Health Status","left"), hdr_cell("Active Alarm Units"),
                   hdr_cell("% of Alarm Units"), hdr_cell("Recommended Action","left")]]

    total_alarm = s["active_alarm_total"] or 1
    data_rows = [(k, s["alarm_by_status"].get(k, 0)) for k in alarm_order if s["alarm_by_status"].get(k, 0) > 0]

    for i, (key, cnt) in enumerate(data_rows):
        bg_c = row_bg(i)
        col = colors3[key]
        pct_str = f"{cnt/total_alarm*100:.1f}%"
        alarm_rows.append([
            data_cell(emojis3[key] + STATUS_CFG.get(key, (key,))[0], "left", True, col, bg_c),
            data_cell(str(cnt), "center", True, col, bg_c),
            data_cell(pct_str, "center", False, C.TEXT, bg_c),
            data_cell(implications.get(key, ""), "left", False, C.SLATE_MID, bg_c),
        ])

    n_data = len(data_rows)
    row_h = [0.40] + [0.44] * n_data
    tbl_h = sum(row_h)
    add_table(slide, 0.38, 2.46, 9.24, tbl_h, alarm_rows,
              col_widths=[2.30, 1.70, 1.70, 3.54],
              row_heights=row_h)

    degraded = s["degraded_with_alarm"]
    total_a = s["active_alarm_total"]
    pct_deg = round(degraded / total_a * 100) if total_a > 0 else 0
    insight_y = 2.46 + tbl_h + 0.12
    add_insight_box(slide, 0.38, insight_y, 9.24, 0.72,
        f"{degraded} out of {total_a} units ({pct_deg}%) with an active VSWR alarm have confirmed RL degradation — "
        f"demonstrating strong alignment between the two indicators. "
        f"These degraded units with active alarms represent the highest-priority field intervention targets in the network.")

    add_text_box(slide, 0.38, 5.42, 9.24, 0.18,
                 "Note: Active VSWR alarm status recorded at time of scan. Units showing no active alarm may have had historical alarm events that have since cleared.",
                 font_size=8, italic=True, color=C.SLATE_MID)


def build_slide_4(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "VSWR Risk Assessment",
                     f"Port return loss measurements — RL degradation threshold: 14 dB (CoG boundary)  |  4 ports per unit  |  {s['vswr_valid_count']:,} valid readings")

    # 2 KPI cards only
    kpi4 = [
        (str(s["risk_units"]),        "VSWR Risk Units",  C.RED_DARK,   C.RED_LIGHT,   C.RED_MID),
        (f"{s['non_risk_units']:,}",  "Non-Risk Units",   C.GREEN_DARK, C.GREEN_LIGHT, C.GREEN_MID),
    ]
    kw4 = 4.50; ky4 = 1.10; kh4 = 0.88
    for i, (val, lbl, dk, lt, nc) in enumerate(kpi4):
        add_kpi_card(slide, 0.38 + i * 4.74, ky4, kw4, kh4, lbl, val, dk, lt, nc)

    # Left: histogram
    add_section_label(slide, 0.38, 2.16, 6.0, "VSWR Distribution (all ports)")
    hist_buf = build_histogram(s["vswr_readings"])
    slide.shapes.add_picture(hist_buf, Inches(0.38), Inches(2.46), Inches(6.0), Inches(2.72))

    # Right: risk summary table
    add_section_label(slide, 6.58, 2.16, 3.04, "Risk Assessment Summary")
    risk_tbl = [
        [hdr_cell("Assessment"), hdr_cell("Units"), hdr_cell("% Fleet")],
        [data_cell("Risk Unit",    "left",   True, C.RED_MID,   C.RED_LIGHT),
         data_cell(str(s["risk_units"]),     "center", True, C.RED_MID,   C.RED_LIGHT),
         data_cell(f"{s['risk_pct']}%",      "center", False, C.TEXT,     C.RED_LIGHT)],
        [data_cell("Non-Risk Unit","left",   True, C.GREEN_MID, C.GREEN_LIGHT),
         data_cell(f"{s['non_risk_units']:,}","center",True, C.GREEN_MID, C.GREEN_LIGHT),
         data_cell(f"{s['non_risk_pct']}%",  "center", False, C.TEXT,     C.GREEN_LIGHT)],
    ]
    add_table(slide, 6.58, 2.46, 3.04, 0.96, risk_tbl,
              col_widths=[1.50, 0.77, 0.77], row_heights=[0.32, 0.32, 0.32])

    # Right: criteria box
    add_rect(slide, 6.58, 3.52, 3.04, 1.58, C.SLATE_LIGHT, C.BORDER, 0.5)
    txb = slide.shapes.add_textbox(Inches(6.70), Inches(3.58), Inches(2.82), Inches(1.48))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run(); r1.text = "Risk Classification Criteria\n"
    r1.font.bold = True; r1.font.size = Pt(10); r1.font.name = "Calibri"; r1.font.color.rgb = C.NAVY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = (
        "A unit is classified as Risk if any port has:\n"
        "  (1) Insufficient downlink power — VSWR measurement unavailable (0 or missing), or\n"
        "  (2) Last measured return loss below 10 dB (VSWR alarm boundary).\n\n"
        f"RL degradation threshold: 14 dB (CoG early warning boundary — 3.1σ below production mean).\n"
        f"Mean fleet VSWR: {s['vswr_mean']} dB."
    )
    r2.font.bold = False; r2.font.size = Pt(9.5); r2.font.name = "Calibri"; r2.font.color.rgb = C.TEXT

    add_text_box(slide, 0.38, 5.20, 6.10, 0.18,
                 f"Key insight: {s['ports_below_14']:,} ports ({s['ports_below_14_pct']}%) fall below 14 dB — {s['ports_below_14_degraded']} from degraded units, "
                 f"{s['ports_below_14_healthy']} from healthy units. Avg {s['avg_degraded_ports']} degraded ports/unit confirms single-port failure mode.",
                 font_size=8.5, color=C.TEXT)
    add_text_box(slide, 0.38, 5.44, 9.24, 0.16,
                 f"Note: 14 dB = CoG early warning boundary in RL algorithm (derived from production CPK: mean=18.6 dB, std=1.47 dB). Risk unit: port NaN/0 (no DL power) or < 10 dB (VSWR alarm boundary). Mean fleet VSWR: {s['vswr_mean']} dB.",
                 font_size=7.5, italic=True, color=C.SLATE_MID)


def build_slide_5(prs, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = C.WHITE

    add_slide_header(slide, "Deployment Age vs Alert Rate",
                     f"RL degradation alert distribution by deployment age cohort  |  units with known deployment date")

    add_section_label(slide, 0.38, 1.12, 9.24, "Alert Unit Breakdown by Deployment Age Cohort")

    def rate_color(r):
        if r > 11: return C.ORANGE_MID
        if r > 9:  return C.AMBER_MID
        return C.TEAL

    def dash(n): return "—" if n == 0 else str(n)

    tbl_rows = [[
        hdr_cell("Age Cohort", "left"),
        hdr_cell("Total (n)"), hdr_cell("Healthy"),
        hdr_cell("Early"), hdr_cell("Middle"),
        hdr_cell("Late"), hdr_cell("Critical"),
        hdr_cell("Alert Total"), hdr_cell("Alert Rate"),
    ]]

    for i, r in enumerate(s["cohort_rows"]):
        bg_c = row_bg(i)
        rc = rate_color(r["rate"])
        coh_label = r["cohort"] + (" *" if r["small"] else "")
        coh_color = C.ORANGE_MID if r["small"] else C.TEXT
        tbl_rows.append([
            data_cell(coh_label,           "left",   False, coh_color,   bg_c),
            data_cell(f"{r['n']:,}",        "center", False, C.TEXT,      bg_c),
            data_cell(f"{r['healthy']:,}",  "center", False, C.GREEN_MID, bg_c),
            data_cell(dash(r["early"]),     "center", False, C.AMBER_MID  if r["early"]  > 0 else C.SLATE_MID, bg_c),
            data_cell(dash(r["middle"]),    "center", False, C.ORANGE_MID if r["middle"] > 0 else C.SLATE_MID, bg_c),
            data_cell(dash(r["late"]),      "center", False, C.RED_MID    if r["late"]   > 0 else C.SLATE_MID, bg_c),
            data_cell(dash(r["critical"]),  "center", False, C.PURPLE_MID if r["critical"]>0 else C.SLATE_MID, bg_c),
            data_cell(str(r["alert"]),      "center", True,  rc,           bg_c),
            data_cell(f"{r['rate']}%",      "center", True,  rc,           bg_c),
        ])

    n_cohorts = len(s["cohort_rows"])
    row_heights = [0.38] + [0.44] * n_cohorts
    add_table(slide, 0.38, 1.42, 9.24, sum(row_heights), tbl_rows,
              col_widths=[1.44, 0.82, 0.90, 0.72, 0.78, 0.72, 0.82, 0.98, 0.94],
              row_heights=row_heights)

    tbl_bottom = 1.42 + sum(row_heights)

    add_text_box(slide, 0.38, tbl_bottom + 0.06, 5.0, 0.20,
                 "* Small sample size — interpret alert rate with caution",
                 font_size=9, italic=True, color=C.ORANGE_MID)

    add_insight_box(slide, 0.38, tbl_bottom + 0.30, 9.24, 0.56,
        "No strong age-degradation correlation observed. Alert rates are broadly uniform across all cohorts, "
        "including recently deployed units. This suggests degradation is driven by installation quality "
        "or environmental factors rather than age-related connector wear.")

    add_text_box(slide, 0.38, 5.42, 9.24, 0.18,
                 "Note: Deployment date from Unit_Field_Deployment_Date field. Units with missing date shown as Unknown.",
                 font_size=8, italic=True, color=C.SLATE_MID)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Generate the Radio 4432 B28 RL Health Analysis Report from monthly CSV scan export."
    )
    parser.add_argument("--input",  "-i", required=True,  help="Path to RRU_BXP_RL50_Health_Result.csv")
    parser.add_argument("--output", "-o", default=None,   help="Output .pptx path (default: RL_Health_Report_<date>.pptx)")
    args = parser.parse_args()

    csv_path = Path(args.input)
    if not csv_path.exists():
        print(f"ERROR: Input file not found: {csv_path}", file=sys.stderr)
        sys.exit(1)

    out_path = args.output or f"RL_Health_Report_{datetime.now().strftime('%Y%m%d')}.pptx"

    print(f"Loading data from: {csv_path}")
    df = load_data(str(csv_path))
    print(f"  {len(df):,} units loaded")

    print("Computing statistics...")
    s = compute_stats(df)
    print(f"  Total units: {s['total_units']:,} | Alert: {s['alert_total']:,} | Healthy: {s['healthy_pct']}%")

    print("Building presentation...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_title(prs, s)
    print("  ✓ Title slide")
    build_slide_1(prs, s)
    print("  ✓ Slide 1 — Executive Summary")
    build_slide_2(prs, s)
    print("  ✓ Slide 2 — Branch Breakdown")
    build_slide_3(prs, s)
    print("  ✓ Slide 3 — Active VSWR Alarm")
    build_slide_4(prs, s)
    print("  ✓ Slide 4 — VSWR Risk Assessment")
    build_slide_5(prs, s)
    print("  ✓ Slide 5 — Deployment Age")

    prs.save(out_path)
    print(f"\nReport saved: {out_path}")


if __name__ == "__main__":
    main()
